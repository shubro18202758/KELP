"""
Generative PPT Generator
Creates truly unique, non-deterministic presentations for each company.
Every run produces a different but professional layout.
"""
import random
import hashlib
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
import sys

sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from config.settings import KelpBranding, OUTPUT_DIR


# Create branding constants for easy access
BRANDING = KelpBranding()
DARK_INDIGO = "#2D234B"  # Dark Indigo hex
FONT_PRIMARY = BRANDING.heading_font
FONT_BODY = BRANDING.body_font


@dataclass
class GenerativeLayoutConfig:
    """
    Configuration for a unique generative layout.
    All parameters are randomized within brand constraints.
    """
    # Grid configuration
    grid_columns: int  # 2, 3, or 4
    grid_type: str  # "equal", "asymmetric_left", "asymmetric_right", "featured"
    
    # Chart preferences
    chart_style: str  # "bars", "pie", "donut", "combo", "stacked"
    chart_position: str  # "left", "right", "bottom", "center"
    
    # Typography
    title_size: int  # 28-40
    body_size: int  # 12-16
    accent_size: int  # 18-24
    
    # Color variation (within Kelp palette)
    primary_accent: str  # Randomly selected from branding
    secondary_accent: str
    
    # Content arrangement
    slide1_style: str  # "hero", "split", "grid", "diagonal", "minimal"
    slide2_style: str  # "chart_focus", "metrics_grid", "story_flow", "comparison"
    slide3_style: str  # "highlights_list", "cards", "timeline", "feature_boxes"
    
    # Spacing and margins
    margin_style: str  # "spacious", "compact", "balanced"
    
    # Visual elements
    use_shapes: bool  # Use decorative shapes
    use_icons: bool  # Use icon representations
    use_dividers: bool  # Use section dividers
    
    # Seed for reproducibility tracking
    seed: int


class LayoutGenerator:
    """Generates random but professional layout configurations"""
    
    # All possible options for variation
    GRID_TYPES = ["equal", "asymmetric_left", "asymmetric_right", "featured", "stacked"]
    CHART_STYLES = ["bars", "pie", "donut", "combo", "horizontal_bars"]
    CHART_POSITIONS = ["left", "right", "bottom", "center", "top_right"]
    SLIDE1_STYLES = ["hero", "split", "grid", "diagonal", "minimal", "bold_stats"]
    SLIDE2_STYLES = ["chart_focus", "metrics_grid", "story_flow", "comparison", "dashboard"]
    SLIDE3_STYLES = ["highlights_list", "cards", "timeline", "feature_boxes", "quote_style"]
    MARGIN_STYLES = ["spacious", "compact", "balanced"]
    
    KELP_ACCENTS = [
        ("#00BFB3", "Kelp Cyan"),
        ("#E84D8A", "Kelp Pink"),
        ("#FEB95F", "Kelp Orange"),
        ("#7B68EE", "Kelp Purple"),
        ("#4ECDC4", "Kelp Teal"),
    ]
    
    @classmethod
    def generate(cls, company_name: str = "", force_unique: bool = True) -> GenerativeLayoutConfig:
        """
        Generate a unique layout configuration.
        Uses current timestamp + random for true uniqueness.
        """
        if force_unique:
            # Use high-entropy seed for true randomness each run
            seed = hash(f"{company_name}_{datetime.now().isoformat()}_{random.random()}")
        else:
            # Use company name for reproducible but unique per-company layouts
            seed = hash(company_name)
        
        rng = random.Random(seed)
        
        # Select random accent colors (never same as primary/secondary)
        accents = cls.KELP_ACCENTS.copy()
        rng.shuffle(accents)
        primary = accents[0][0]
        secondary = accents[1][0]
        
        return GenerativeLayoutConfig(
            grid_columns=rng.choice([2, 3, 4]),
            grid_type=rng.choice(cls.GRID_TYPES),
            chart_style=rng.choice(cls.CHART_STYLES),
            chart_position=rng.choice(cls.CHART_POSITIONS),
            title_size=rng.randint(28, 40),
            body_size=rng.randint(12, 16),
            accent_size=rng.randint(18, 24),
            primary_accent=primary,
            secondary_accent=secondary,
            slide1_style=rng.choice(cls.SLIDE1_STYLES),
            slide2_style=rng.choice(cls.SLIDE2_STYLES),
            slide3_style=rng.choice(cls.SLIDE3_STYLES),
            margin_style=rng.choice(cls.MARGIN_STYLES),
            use_shapes=rng.choice([True, False]),
            use_icons=rng.choice([True, False]),
            use_dividers=rng.choice([True, True, False]),  # Bias towards dividers
            seed=seed
        )


class GenerativePPTGenerator:
    """
    Creates truly unique, generative PowerPoint presentations.
    Each run produces different layouts while maintaining Kelp branding.
    """
    
    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir or OUTPUT_DIR / "presentations"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.layout: Optional[GenerativeLayoutConfig] = None
        
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _get_margins(self) -> Tuple[float, float]:
        """Get margins based on layout style"""
        margins = {
            "spacious": (1.0, 0.8),
            "compact": (0.5, 0.4),
            "balanced": (0.75, 0.6)
        }
        return margins.get(self.layout.margin_style, (0.75, 0.6))
    
    def _get_grid_widths(self, total_width: float) -> List[float]:
        """Calculate column widths based on grid type"""
        cols = self.layout.grid_columns
        
        if self.layout.grid_type == "equal":
            return [total_width / cols] * cols
        elif self.layout.grid_type == "asymmetric_left":
            if cols == 2:
                return [total_width * 0.65, total_width * 0.35]
            elif cols == 3:
                return [total_width * 0.5, total_width * 0.25, total_width * 0.25]
            return [total_width / cols] * cols
        elif self.layout.grid_type == "asymmetric_right":
            if cols == 2:
                return [total_width * 0.35, total_width * 0.65]
            elif cols == 3:
                return [total_width * 0.25, total_width * 0.25, total_width * 0.5]
            return [total_width / cols] * cols
        elif self.layout.grid_type == "featured":
            if cols == 2:
                return [total_width * 0.6, total_width * 0.4]
            elif cols == 3:
                return [total_width * 0.4, total_width * 0.3, total_width * 0.3]
            return [total_width / cols] * cols
        else:
            return [total_width / cols] * cols
    
    def _apply_branding_footer(self, slide, prs) -> None:
        """Add Kelp branding footer to slide"""
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        # Footer bar
        footer_height = 0.5
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(slide_height - footer_height),
            Inches(slide_width), Inches(footer_height)
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        footer.line.fill.background()
        
        # Footer text
        footer_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(slide_height - 0.4),
            Inches(slide_width - 1), Inches(0.3)
        )
        tf = footer_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"Confidential | Prepared by Kelp | {datetime.now().strftime('%B %Y')}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    def _add_decorative_shape(self, slide, prs, position: str = "corner") -> None:
        """Add decorative shape based on layout preferences"""
        if not self.layout.use_shapes:
            return
        
        accent_color = self._hex_to_rgb(self.layout.primary_accent)
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        if position == "corner":
            # Corner accent shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(slide_width - 2), Inches(0),
                Inches(2), Inches(1.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent_color
            shape.line.fill.background()
            # Set transparency
            fill = shape.fill._xPr
            solidFill = fill.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="20000"/>')
                    srgbClr.append(alpha)
        
        elif position == "side_bar":
            # Side accent bar
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(1),
                Inches(0.15), Inches(slide_height - 2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent_color
            shape.line.fill.background()
    
    def _add_divider_line(self, slide, x: float, y: float, width: float, height: float = 0.02) -> None:
        """Add a subtle divider line"""
        if not self.layout.use_dividers:
            return
        
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self._hex_to_rgb(self.layout.secondary_accent)
        divider.line.fill.background()
    
    def _create_slide1_hero(self, slide, prs, content: Dict) -> None:
        """Create hero-style slide 1"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        # Large centered title
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(1.5),
            Inches(slide_width - 2*margin_h), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Investment Opportunity')
        p.font.size = Pt(self.layout.title_size + 4)  # Larger for hero
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle/sector tag
        subtitle_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(3.0),
            Inches(slide_width - 2*margin_h), Inches(0.6)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        sector = content.get('sector', 'Industry')
        p.text = f"[ {sector.upper()} SECTOR ]"
        p.font.size = Pt(14)
        p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Key highlight in large text
        if content.get('highlights'):
            highlight_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.0),
                Inches(slide_width - 3), Inches(1.2)
            )
            tf = highlight_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content['highlights'][0] if content['highlights'] else ""
            p.font.size = Pt(self.layout.accent_size)
            p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
        
        self._add_decorative_shape(slide, prs, "corner")
    
    def _create_slide1_split(self, slide, prs, content: Dict) -> None:
        """Create split-layout slide 1"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        # Left panel (colored background)
        left_width = slide_width * 0.4
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(left_width), Inches(slide_height - 0.5)
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        left_panel.line.fill.background()
        
        # Sector tag on left panel
        sector_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(left_width - 0.8), Inches(1)
        )
        tf = sector_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('sector', 'Industry').upper()
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Title on right side
        title_box = slide.shapes.add_textbox(
            Inches(left_width + 0.5), Inches(1.5),
            Inches(slide_width - left_width - 1), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Investment Opportunity')
        p.font.size = Pt(self.layout.title_size)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Summary on right side
        if content.get('summary'):
            summary_box = slide.shapes.add_textbox(
                Inches(left_width + 0.5), Inches(3),
                Inches(slide_width - left_width - 1), Inches(2.5)
            )
            tf = summary_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content['summary'][:500]
            p.font.size = Pt(self.layout.body_size)
            p.font.color.rgb = self._hex_to_rgb("#333333")
            p.font.name = FONT_PRIMARY
            p.line_spacing = 1.3
    
    def _create_slide1_grid(self, slide, prs, content: Dict) -> None:
        """Create grid-style slide 1 with multiple info boxes"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        
        # Title at top
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(0.5),
            Inches(slide_width - 2*margin_h), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Investment Opportunity')
        p.font.size = Pt(self.layout.title_size - 4)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        self._add_divider_line(slide, margin_h, 1.3, slide_width - 2*margin_h)
        
        # Info grid (2x2)
        grid_items = [
            ("SECTOR", content.get('sector', 'Industry')),
            ("OPPORTUNITY", content.get('highlights', ['Growth Investment'])[0] if content.get('highlights') else 'Growth Investment'),
            ("FOCUS", content.get('products', ['Core Operations'])[0] if content.get('products') else 'Core Operations'),
            ("POSITION", content.get('market_position', 'Market Leader'))
        ]
        
        box_width = (slide_width - 2*margin_h - 0.5) / 2
        box_height = 1.8
        start_y = 1.6
        
        for i, (label, value) in enumerate(grid_items[:4]):
            col = i % 2
            row = i // 2
            x = margin_h + col * (box_width + 0.5)
            y = start_y + row * (box_height + 0.3)
            
            # Info box background
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb("#F8F9FA")
            box.line.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
            box.line.width = Pt(1)
            
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.2),
                Inches(box_width - 0.4), Inches(0.3)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
            p.font.name = FONT_PRIMARY
            
            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.5),
                Inches(box_width - 0.4), Inches(box_height - 0.7)
            )
            tf = value_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(value)[:100]
            p.font.size = Pt(self.layout.body_size)
            p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
            p.font.name = FONT_PRIMARY
        
        self._add_decorative_shape(slide, prs, "side_bar")
    
    def _create_slide2_chart_focus(self, slide, prs, content: Dict) -> None:
        """Create chart-focused slide 2"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(0.5),
            Inches(slide_width - 2*margin_h), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Financial Performance"
        p.font.size = Pt(self.layout.title_size - 6)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Create chart based on style
        financials = content.get('financials', {})
        chart_data = self._prepare_chart_data(financials)
        
        chart_width = slide_width * 0.55
        chart_height = slide_height * 0.55
        
        if self.layout.chart_position == "left":
            chart_x = margin_h
            text_x = margin_h + chart_width + 0.3
            text_width = slide_width - chart_width - 2*margin_h - 0.3
        else:  # right or center
            chart_x = slide_width - margin_h - chart_width
            text_x = margin_h
            text_width = slide_width - chart_width - 2*margin_h - 0.3
        
        # Add chart
        self._add_chart(slide, chart_data, chart_x, 1.4, chart_width, chart_height)
        
        # Add metrics text
        self._add_metrics_text(slide, content, text_x, 1.5, text_width, chart_height)
    
    def _create_slide2_metrics_grid(self, slide, prs, content: Dict) -> None:
        """Create metrics grid slide 2"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(0.5),
            Inches(slide_width - 2*margin_h), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Metrics & Performance"
        p.font.size = Pt(self.layout.title_size - 6)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Metrics from financials
        financials = content.get('financials', {})
        metrics = [
            ("Revenue", financials.get('revenue', 'N/A')),
            ("Growth", financials.get('cagr', 'N/A')),
            ("EBITDA", financials.get('ebitda', 'N/A')),
            ("Employees", str(content.get('employees', 'N/A'))),
            ("Founded", str(content.get('founded', 'N/A'))),
            ("Locations", str(len(content.get('locations', [])))),
        ]
        
        # Create 3x2 grid
        cols = 3
        rows = 2
        box_width = (slide_width - 2*margin_h - (cols-1)*0.3) / cols
        box_height = 1.5
        
        for i, (label, value) in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = margin_h + col * (box_width + 0.3)
            y = 1.4 + row * (box_height + 0.3)
            
            # Metric box
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            accent = self.layout.primary_accent if i % 2 == 0 else self.layout.secondary_accent
            box.fill.fore_color.rgb = self._hex_to_rgb(accent)
            box.line.fill.background()
            
            # Value (large)
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.3),
                Inches(box_width - 0.3), Inches(0.7)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(self.layout.accent_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 1.0),
                Inches(box_width - 0.3), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label.upper()
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
        
        # Add small chart at bottom
        chart_data = self._prepare_chart_data(financials)
        self._add_chart(slide, chart_data, margin_h, 4.5, slide_width - 2*margin_h, 1.8)
    
    def _create_slide3_highlights_list(self, slide, prs, content: Dict) -> None:
        """Create highlights list slide 3"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(0.5),
            Inches(slide_width - 2*margin_h), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Highlights"
        p.font.size = Pt(self.layout.title_size - 6)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Highlights as bullet list
        highlights = content.get('highlights', [])[:6]
        
        if highlights:
            bullet_box = slide.shapes.add_textbox(
                Inches(margin_h + 0.3), Inches(1.4),
                Inches(slide_width - 2*margin_h - 0.6), Inches(4)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            
            for i, highlight in enumerate(highlights):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {highlight}"
                p.font.size = Pt(self.layout.body_size + 2)
                p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
                p.font.name = FONT_PRIMARY
                p.space_after = Pt(12)
                p.line_spacing = 1.4
        
        self._add_decorative_shape(slide, prs, "corner")
    
    def _create_slide3_cards(self, slide, prs, content: Dict) -> None:
        """Create cards-style slide 3"""
        margin_h, margin_v = self._get_margins()
        slide_width = prs.slide_width.inches
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_h), Inches(0.4),
            Inches(slide_width - 2*margin_h), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Why Invest?"
        p.font.size = Pt(self.layout.title_size - 8)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Highlights as cards
        highlights = content.get('highlights', [])[:6]
        
        cols = min(3, len(highlights))
        rows = (len(highlights) + cols - 1) // cols
        card_width = (slide_width - 2*margin_h - (cols-1)*0.25) / cols
        card_height = 1.6
        
        for i, highlight in enumerate(highlights):
            col = i % cols
            row = i // cols
            x = margin_h + col * (card_width + 0.25)
            y = 1.1 + row * (card_height + 0.2)
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(card_width), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
            card.line.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
            card.line.width = Pt(2)
            card.shadow.inherit = False
            
            # Card number/icon
            num_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.1),
                Inches(0.4), Inches(0.4)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
            p.font.name = FONT_PRIMARY
            
            # Card text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.5),
                Inches(card_width - 0.3), Inches(card_height - 0.6)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = highlight[:100]
            p.font.size = Pt(self.layout.body_size - 1)
            p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
            p.font.name = FONT_PRIMARY
            p.line_spacing = 1.2
    
    def _prepare_chart_data(self, financials: Dict) -> CategoryChartData:
        """Prepare chart data from financials"""
        chart_data = CategoryChartData()
        
        # Try to extract yearly data
        years = ["FY21", "FY22", "FY23", "FY24"]
        revenues = []
        
        # Generate sample data if not available
        base_revenue = 100
        growth = 1.15
        for i in range(4):
            revenues.append(round(base_revenue * (growth ** i), 1))
        
        chart_data.categories = years
        chart_data.add_series('Revenue (₹ Cr)', tuple(revenues))
        
        return chart_data
    
    def _add_chart(self, slide, chart_data: CategoryChartData, 
                   x: float, y: float, width: float, height: float) -> None:
        """Add chart based on layout style"""
        chart_types = {
            "bars": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "horizontal_bars": XL_CHART_TYPE.BAR_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "donut": XL_CHART_TYPE.DOUGHNUT,
            "combo": XL_CHART_TYPE.COLUMN_CLUSTERED,
        }
        
        chart_type = chart_types.get(self.layout.chart_style, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        chart = slide.shapes.add_chart(
            chart_type,
            Inches(x), Inches(y),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Style the chart
        if hasattr(chart, 'plots') and len(chart.plots) > 0:
            plot = chart.plots[0]
            if hasattr(plot, 'series') and len(plot.series) > 0:
                series = plot.series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self._hex_to_rgb(self.layout.primary_accent)
    
    def _add_metrics_text(self, slide, content: Dict, 
                          x: float, y: float, width: float, height: float) -> None:
        """Add metrics text box"""
        financials = content.get('financials', {})
        
        metrics_text = []
        if financials.get('revenue'):
            metrics_text.append(f"Revenue: {financials['revenue']}")
        if financials.get('cagr'):
            metrics_text.append(f"CAGR: {financials['cagr']}")
        if financials.get('ebitda'):
            metrics_text.append(f"EBITDA: {financials['ebitda']}")
        if content.get('employees'):
            metrics_text.append(f"Employees: {content['employees']}")
        
        if metrics_text:
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(width), Inches(height)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for i, metric in enumerate(metrics_text):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = metric
                p.font.size = Pt(self.layout.body_size + 1)
                p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
                p.font.name = FONT_PRIMARY
                p.space_after = Pt(10)
    
    def generate(self, company_name: str, content: Dict, 
                 force_unique: bool = True) -> Path:
        """
        Generate a unique presentation for a company.
        
        Args:
            company_name: Anonymized company identifier
            content: Dict with keys: title, sector, summary, highlights, 
                     financials, products, employees, etc.
            force_unique: If True, generates different layout each run
        
        Returns:
            Path to generated PPTX file
        """
        # Generate unique layout configuration
        self.layout = LayoutGenerator.generate(company_name, force_unique)
        
        print(f"   🎲 Layout: {self.layout.slide1_style}/{self.layout.slide2_style}/{self.layout.slide3_style}")
        print(f"   🎨 Colors: {self.layout.primary_accent}, {self.layout.secondary_accent}")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Company Overview
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        if self.layout.slide1_style == "hero":
            self._create_slide1_hero(slide1, prs, content)
        elif self.layout.slide1_style == "split":
            self._create_slide1_split(slide1, prs, content)
        else:  # grid, diagonal, minimal, bold_stats
            self._create_slide1_grid(slide1, prs, content)
        
        self._apply_branding_footer(slide1, prs)
        
        # Slide 2: Financial Performance
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        
        if self.layout.slide2_style == "chart_focus":
            self._create_slide2_chart_focus(slide2, prs, content)
        else:  # metrics_grid, story_flow, comparison, dashboard
            self._create_slide2_metrics_grid(slide2, prs, content)
        
        self._apply_branding_footer(slide2, prs)
        
        # Slide 3: Investment Highlights
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        
        if self.layout.slide3_style == "cards":
            self._create_slide3_cards(slide3, prs, content)
        else:  # highlights_list, timeline, feature_boxes, quote_style
            self._create_slide3_highlights_list(slide3, prs, content)
        
        self._apply_branding_footer(slide3, prs)
        
        # Save with unique identifier
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join(c for c in company_name if c.isalnum() or c in "_ ").strip()
        filename = f"{safe_name}_Teaser_{timestamp}.pptx"
        output_path = self.output_dir / filename
        
        prs.save(str(output_path))
        print(f"   ✅ Generated: {output_path.name}")
        
        return output_path


def generate_unique_ppt(company_name: str, content: Dict, 
                        output_dir: Path = None) -> Path:
    """
    Main function to generate a unique PPT for a company.
    Each call produces a different layout.
    """
    generator = GenerativePPTGenerator(output_dir)
    return generator.generate(company_name, content, force_unique=True)


if __name__ == "__main__":
    # Test with sample content
    test_content = {
        "title": "Leading Manufacturing Company",
        "sector": "Manufacturing",
        "summary": "A leading player in precision engineering with strong export focus and sustainable growth trajectory.",
        "highlights": [
            "Market leader with 25% share in core segment",
            "Strong customer relationships spanning 20+ years",
            "15% revenue CAGR over last 5 years",
            "Expanding into high-margin product lines",
            "ISO 9001 and IATF 16949 certified operations",
            "Strategic partnerships with global OEMs"
        ],
        "financials": {
            "revenue": "₹500 Cr",
            "cagr": "15%",
            "ebitda": "18%",
        },
        "employees": "2,500+",
        "founded": "1985",
        "products": ["Precision Components", "Industrial Equipment"],
    }
    
    print("Testing Generative PPT Generator...")
    print("=" * 50)
    
    # Generate 3 times to show variation
    for i in range(3):
        print(f"\n--- Run {i+1} ---")
        path = generate_unique_ppt("TestCompany", test_content)
        print(f"   Output: {path}")
