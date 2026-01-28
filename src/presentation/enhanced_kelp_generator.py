"""
Enhanced Kelp Investment Teaser Generator V2
=============================================

Creates data-dense, visually rich investment presentations with:
- Multiple charts and graphs per slide
- Web-sourced sector images
- Efficient space utilization
- Professional M&A formatting

Slide Structure:
- Slide 1: Cover with sector imagery
- Slide 2: Business Overview (dense 4-quadrant layout)
- Slide 3: Financial Deep-Dive (3+ charts, metrics grid)
- Slide 4: Investment Thesis (5 highlights with supporting data)
- Slide 5: Disclaimer
"""

import random
import math
import asyncio
import aiohttp
import io
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import nsmap
import sys

# Try to import PIL for image handling
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from config.settings import OUTPUT_DIR


# ============================================================================
# KELP BRANDING - Enhanced Color Palette
# ============================================================================

class KelpBrandV2:
    """Enhanced Kelp brand with more colors for rich visuals"""
    # Core brand colors
    NAVY = "#1B365D"
    DARK_BLUE = "#0D2137"
    TEAL = "#00A5B5"
    CYAN = "#00BFB3"
    PINK = "#E84B8A"
    ORANGE = "#F5A623"
    PURPLE = "#7B68EE"
    
    # Chart colors (matching reference)
    CHART_NAVY = "#2B5797"
    CHART_ORANGE = "#E98B39"
    CHART_PURPLE = "#7030A0"
    CHART_GREEN = "#4CAF50"
    CHART_RED = "#E53935"
    CHART_TEAL = "#00897B"
    
    # Neutrals
    WHITE = "#FFFFFF"
    LIGHT_GRAY = "#F5F5F5"
    MED_GRAY = "#E0E0E0"
    DARK_GRAY = "#333333"
    MID_GRAY = "#666666"
    
    # Accent backgrounds
    LIGHT_BLUE = "#E3F2FD"
    LIGHT_ORANGE = "#FFF3E0"
    LIGHT_GREEN = "#E8F5E9"
    LIGHT_PURPLE = "#EDE7F6"
    LIGHT_YELLOW = "#FFFDE7"
    
    # Typography
    FONT_PRIMARY = "Arial"
    FONT_HEADING = "Arial"
    
    # Slide size 16:9
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)


# Project codenames
PROJECT_CODENAMES = [
    "Alpha", "Apex", "Atlas", "Aurora", "Catalyst", "Cosmos", "Eclipse", 
    "Everest", "Falcon", "Frontier", "Galaxy", "Helios", "Horizon", 
    "Jupiter", "Mercury", "Momentum", "Neptune", "Nexus", "Nova", "Omega",
    "Orion", "Phoenix", "Pinnacle", "Prime", "Quantum", "Saturn", "Stellar",
    "Stratos", "Summit", "Titan", "Velocity", "Vertex", "Zenith"
]


@dataclass
class EnhancedTeaserData:
    """Complete data for enhanced teaser generation"""
    # Core
    codename: str = ""
    sector: str = "Manufacturing"
    sub_sector: str = ""
    
    # Business Overview (Slide 2) - 4 Quadrants
    business_bullets: List[str] = field(default_factory=list)
    key_customers: List[str] = field(default_factory=list)
    at_a_glance: Dict[str, str] = field(default_factory=dict)
    product_portfolio: List[str] = field(default_factory=list)
    applications: List[str] = field(default_factory=list)
    certifications: List[str] = field(default_factory=list)
    
    # Financial Deep-Dive (Slide 3) - Multiple Charts
    key_metrics_headline: Dict[str, str] = field(default_factory=dict)
    revenue_trend: Dict[str, List] = field(default_factory=dict)
    margin_trend: Dict[str, List] = field(default_factory=dict)
    segment_mix: Dict[str, float] = field(default_factory=dict)
    geographic_mix: Dict[str, float] = field(default_factory=dict)
    growth_callouts: List[str] = field(default_factory=list)
    capex_plans: List[str] = field(default_factory=list)
    
    # Investment Highlights (Slide 4)
    investment_highlights: List[Dict[str, str]] = field(default_factory=list)
    thesis_points: List[str] = field(default_factory=list)  # Investment thesis bullets
    key_risks: List[str] = field(default_factory=list)
    
    # Market Intelligence (from Web Research)
    market_size: str = ""  # e.g., "$45 Billion by 2027"
    market_cagr: str = ""  # e.g., "12.5% CAGR"
    industry_trends: List[str] = field(default_factory=list)
    competitive_landscape: str = ""  # Overview of competitors
    
    # Images
    sector_images: List[Path] = field(default_factory=list)
    chart_images: List[Path] = field(default_factory=list)
    
    # Sources
    data_sources: List[str] = field(default_factory=list)


class EnhancedKelpGenerator:
    """
    Creates data-dense, visually rich investment teasers.
    
    Key Improvements:
    - 4-quadrant business overview layout
    - Multiple charts on financial slide
    - Metrics grids and KPI dashboards
    - Better space utilization
    - Icon-based visual elements
    - Sector-appropriate images on every slide
    """
    
    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir or OUTPUT_DIR / "enhanced_teasers"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs: Optional[Presentation] = None
        self.brand = KelpBrandV2()
        
        # Image storage per slide
        self.slide_images: Dict[str, List[Path]] = {}
    
    def _add_image_to_slide(self, slide, image_path: Path, x: float, y: float,
                           width: float = None, height: float = None,
                           opacity: float = 1.0, send_to_back: bool = False) -> bool:
        """
        Add an image to a slide at the specified position.
        
        Args:
            slide: PowerPoint slide object
            image_path: Path to image file
            x, y: Position in inches
            width, height: Size in inches (if None, uses image aspect ratio)
            opacity: Image opacity (0.0 to 1.0) - placeholder for future use
            send_to_back: If True, send image behind other elements
        
        Returns:
            True if successful, False otherwise
        """
        if not image_path or not image_path.exists():
            return False
        
        try:
            # Open image to get dimensions
            if HAS_PIL:
                with Image.open(image_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    # Calculate dimensions
                    if width and not height:
                        height = width / aspect_ratio
                    elif height and not width:
                        width = height * aspect_ratio
                    elif not width and not height:
                        width = 2.0  # Default width
                        height = width / aspect_ratio
            else:
                if not width:
                    width = 2.0
                if not height:
                    height = 1.5
            
            # Add image to slide
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(x), Inches(y),
                Inches(width), Inches(height)
            )
            
            # Send to back if requested
            if send_to_back:
                spTree = slide.shapes._spTree
                sp = picture._element
                spTree.remove(sp)
                spTree.insert(3, sp)  # Insert after background
            
            return True
            
        except Exception as e:
            print(f"    ⚠ Failed to add image {image_path.name}: {e}")
            return False
    
    def set_slide_images(self, slide_images: Dict[str, List[Path]]) -> None:
        """Set images for all slides before generation"""
        self.slide_images = slide_images or {}
        
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    
    def _add_kelp_logo(self, slide, x: float, y: float, dark_bg: bool = True) -> None:
        """Add Kelp logo"""
        logo = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.5))
        tf = logo.text_frame
        p = tf.paragraphs[0]
        
        run1 = p.add_run()
        run1.text = "* "
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = self._hex_to_rgb(KelpBrandV2.PINK)
        
        run2 = p.add_run()
        run2.text = "Kelp"
        run2.font.size = Pt(18)
        run2.font.bold = True
        run2.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE if dark_bg else KelpBrandV2.NAVY)
        run2.font.name = KelpBrandV2.FONT_HEADING
    
    def _add_footer(self, slide) -> None:
        """Add standard footer"""
        footer = slide.shapes.add_textbox(Inches(0), Inches(7.15), Inches(13.333), Inches(0.3))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
        p.font.size = Pt(8)
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.MID_GRAY)
        p.font.name = KelpBrandV2.FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    def _add_header_bar(self, slide, title: str) -> None:
        """Add header bar with title and logo"""
        # Header background
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                        KelpBrandV2.SLIDE_WIDTH, Inches(0.55))
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(8), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.font.name = KelpBrandV2.FONT_HEADING
        
        # Logo
        self._add_kelp_logo(slide, 11.8, 0.05, dark_bg=True)
    
    def _add_metric_card(self, slide, x: float, y: float, w: float, h: float,
                         value: str, label: str, color: str, icon: str = "") -> None:
        """Add a metric card with value and label"""
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                      Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = self._hex_to_rgb(color)
        card.line.fill.background()
        
        # Small decorative dot instead of emoji icon (renders reliably)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                     Inches(x + 0.08), Inches(y + 0.08),
                                     Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        dot.fill.fore_color.brightness = 0.3  # Slightly transparent
        dot.line.fill.background()
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08),
                                           Inches(w - 0.2), Inches(h * 0.55))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + h * 0.55),
                                           Inches(w - 0.1), Inches(h * 0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label.upper()[:15]
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_section_box(self, slide, x: float, y: float, w: float, h: float,
                         title: str, items: List[str], header_color: str,
                         bg_color: str = None) -> None:
        """Add a section box with header and bullet items"""
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                    Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(bg_color or KelpBrandV2.WHITE)
        bg.line.color.rgb = self._hex_to_rgb(KelpBrandV2.MED_GRAY)
        bg.line.width = Pt(0.75)
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                        Inches(w), Inches(0.35))
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(header_color)
        header.line.fill.background()
        
        # Header text
        header_txt = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05),
                                              Inches(w - 0.2), Inches(0.3))
        tf = header_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Items
        if items:
            item_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.42),
                                                Inches(w - 0.2), Inches(h - 0.5))
            tf = item_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(items[:8]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item[:60]}"
                p.font.size = Pt(9)
                p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
                p.space_after = Pt(3)
    
    def _add_chart(self, slide, x: float, y: float, w: float, h: float,
                   chart_type: str, data: Dict, title: str = "",
                   colors: List[str] = None) -> None:
        """
        Add a chart with extensive customization.
        
        Supported chart types:
        - bar, column: Clustered bar/column charts
        - stacked_column, stacked_bar: Stacked variants
        - line, line_markers: Line charts with/without markers
        - area, stacked_area: Area charts
        - pie, donut: Circular charts with data labels
        - combo: Combination bar + line chart
        - waterfall: Waterfall/bridge chart (simulated)
        """
        chart_data = CategoryChartData()
        
        if chart_type in ['bar', 'column', 'line', 'line_markers', 'area', 
                          'stacked_column', 'stacked_bar', 'stacked_area', 'combo']:
            categories = data.get('categories', ['A', 'B', 'C', 'D'])
            chart_data.categories = categories
            
            # Support multiple series
            series_data = data.get('series', [data.get('values', [100, 120, 140, 160])])
            series_names = data.get('series_names', ['Value'])
            
            if isinstance(series_data[0], (list, tuple)):
                for i, series in enumerate(series_data):
                    name = series_names[i] if i < len(series_names) else f'Series {i+1}'
                    chart_data.add_series(name, tuple(series))
            else:
                chart_data.add_series(series_names[0], tuple(series_data))
        
        elif chart_type in ['pie', 'donut']:
            values = data.get('values', {"A": 60, "B": 40})
            if isinstance(values, dict):
                categories = data.get('labels', list(values.keys()))
                values = list(values.values())
            else:
                categories = data.get('labels', ['Segment A', 'Segment B'])
            chart_data.categories = categories
            series_name = title if title else 'Revenue %'
            chart_data.add_series(series_name, tuple(values))
        
        # Extended chart type mapping
        xl_types = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
            'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
            'line': XL_CHART_TYPE.LINE,
            'line_markers': XL_CHART_TYPE.LINE_MARKERS,
            'area': XL_CHART_TYPE.AREA,
            'stacked_area': XL_CHART_TYPE.AREA_STACKED,
            'pie': XL_CHART_TYPE.PIE,
            'donut': XL_CHART_TYPE.DOUGHNUT,
            'combo': XL_CHART_TYPE.COLUMN_CLUSTERED,  # Will add line overlay
        }
        
        chart = slide.shapes.add_chart(
            xl_types.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
            Inches(x), Inches(y), Inches(w), Inches(h), chart_data
        )
        
        chart_obj = chart.chart
        
        # Style the chart
        if chart_obj.has_legend:
            chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart_obj.legend.include_in_layout = False
            chart_obj.legend.font.size = Pt(8)
        
        # Add data labels for pie/donut charts
        if chart_type in ['pie', 'donut']:
            try:
                plot = chart_obj.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.show_category_name = True
                data_labels.show_percentage = True
                data_labels.show_value = False
                data_labels.font.size = Pt(9)
                data_labels.font.bold = True
            except Exception:
                pass
        
        # Add data labels for column/bar charts if show_labels requested
        if chart_type in ['column', 'bar', 'stacked_column'] and data.get('show_labels', False):
            try:
                for series in chart_obj.series:
                    series.has_data_labels = True
                    series.data_labels.font.size = Pt(8)
                    series.data_labels.font.bold = True
            except Exception:
                pass
        
        # Apply colors
        if colors:
            for i, series in enumerate(chart_obj.series):
                if i < len(colors):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = self._hex_to_rgb(colors[i])
        
        # Add title if provided
        if title:
            title_box = slide.shapes.add_textbox(Inches(x), Inches(y - 0.25),
                                                 Inches(w), Inches(0.25))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_kpi_gauge(self, slide, x: float, y: float, value: str, label: str,
                       color: str = None, size: float = 1.0) -> None:
        """Add a circular KPI gauge/indicator"""
        color = color or KelpBrandV2.TEAL
        
        # Outer ring
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                      Inches(x), Inches(y),
                                      Inches(size), Inches(size))
        ring.fill.solid()
        ring.fill.fore_color.rgb = self._hex_to_rgb(color)
        ring.line.fill.background()
        
        # Inner circle (white)
        inner = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(x + size * 0.15), Inches(y + size * 0.15),
                                       Inches(size * 0.7), Inches(size * 0.7))
        inner.fill.solid()
        inner.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        inner.line.fill.background()
        
        # Value text
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + size * 0.3),
                                           Inches(size), Inches(size * 0.4))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(int(14 * size))
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(color)
        p.alignment = PP_ALIGN.CENTER
        
        # Label below
        lbl_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(y + size + 0.05),
                                           Inches(size + 0.4), Inches(0.25))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_trend_indicator(self, slide, x: float, y: float, 
                             value: str, change: str, label: str,
                             is_positive: bool = True) -> None:
        """Add a trend indicator with arrow"""
        # Background card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(y),
                                      Inches(1.5), Inches(0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = self._hex_to_rgb(
            KelpBrandV2.LIGHT_GREEN if is_positive else KelpBrandV2.LIGHT_ORANGE
        )
        card.line.fill.background()
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1),
                                           Inches(1.0), Inches(0.35))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        
        # Change indicator
        arrow = "↑" if is_positive else "↓"
        change_color = KelpBrandV2.CHART_GREEN if is_positive else KelpBrandV2.ORANGE
        
        chg_box = slide.shapes.add_textbox(Inches(x + 1.0), Inches(y + 0.12),
                                           Inches(0.4), Inches(0.3))
        tf = chg_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{arrow}{change}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(change_color)
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.48),
                                           Inches(1.3), Inches(0.25))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.MID_GRAY)
    
    # ========================================================================
    # SLIDE 1: COVER
    # ========================================================================
    
    def _create_slide1_cover(self, slide, data: EnhancedTeaserData) -> None:
        """Create cover slide with gradient and sector imagery"""
        # Gradient background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    KelpBrandV2.SLIDE_WIDTH, KelpBrandV2.SLIDE_HEIGHT)
        bg.fill.gradient()
        bg.fill.gradient_angle = 135
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_BLUE)
        stops[1].color.rgb = self._hex_to_rgb(KelpBrandV2.PURPLE)
        bg.line.fill.background()
        
        # Left accent bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                        Inches(0.35), KelpBrandV2.SLIDE_HEIGHT)
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        accent.line.fill.background()
        
        # Geometric accent (top right)
        tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8), Inches(-0.5),
                                     Inches(6), Inches(4.5))
        tri.fill.solid()
        tri.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.PINK)
        tri.fill.fore_color.brightness = 0.15
        tri.line.fill.background()
        
        # Send backgrounds to back
        for shape in [bg, accent, tri]:
            spTree = slide.shapes._spTree
            sp = shape._element
            spTree.remove(sp)
            spTree.insert(2, sp)
        
        # Logo
        self._add_kelp_logo(slide, 0.7, 0.6)
        
        # Project codename
        title = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(8), Inches(1.2))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = f"Project {data.codename}"
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Sector tag
        sector_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                            Inches(0.7), Inches(3.6), Inches(3.5), Inches(0.4))
        sector_box.fill.solid()
        sector_box.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.TEAL)
        sector_box.line.fill.background()
        
        sector_txt = slide.shapes.add_textbox(Inches(0.7), Inches(3.65), Inches(3.5), Inches(0.35))
        tf = sector_txt.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {data.sector}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(5), Inches(0.8))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = "Confidential Investment Brief"
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Date and website
        info = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(4), Inches(0.6))
        tf = info.text_frame
        p = tf.paragraphs[0]
        p.text = f"January 2026 | kelpglobal.com"
        p.font.size = Pt(11)
        p.font.color.rgb = self._hex_to_rgb("#AAAAAA")
        
        # ===== ADD SECTOR IMAGES (Right side - hero image area) =====
        slide1_images = self.slide_images.get('slide1', [])
        if slide1_images:
            # Main hero image (right side, prominent position)
            if len(slide1_images) >= 1:
                img_path = slide1_images[0].path if hasattr(slide1_images[0], 'path') else slide1_images[0]
                self._add_image_to_slide(slide, img_path, 7.5, 1.5, width=5.5, height=4.0, send_to_back=False)
            # Smaller accent image (top right corner)
            if len(slide1_images) >= 2:
                img_path = slide1_images[1].path if hasattr(slide1_images[1], 'path') else slide1_images[1]
                self._add_image_to_slide(slide, img_path, 10.5, 5.7, width=2.5, height=1.6, send_to_back=False)
    
    # ========================================================================
    # SLIDE 2: BUSINESS OVERVIEW (Redesigned - More Visual & Engaging)
    # ========================================================================
    
    def _create_slide2_business_overview(self, slide, data: EnhancedTeaserData) -> None:
        """Create data-rich business overview with proper proportions"""
        # Clean white background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    KelpBrandV2.SLIDE_WIDTH, KelpBrandV2.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header
        self._add_header_bar(slide, "Business Overview")
        
        # ===== TOP ROW: Image + Key Metrics (side by side) =====
        slide2_images = self.slide_images.get('slide2', [])
        
        # Hero image - 16:9 aspect ratio (3.2 x 1.8)
        if slide2_images and len(slide2_images) >= 1:
            img_path = slide2_images[0].path if hasattr(slide2_images[0], 'path') else slide2_images[0]
            self._add_image_to_slide(slide, img_path, 0.2, 0.7, width=4.8, height=2.7, send_to_back=False)
        
        # Key Metrics - 4 large cards in 2x2 grid
        metrics = list(data.at_a_glance.items())[:4]
        metric_colors = [KelpBrandV2.ORANGE, KelpBrandV2.TEAL, KelpBrandV2.PURPLE, KelpBrandV2.CHART_GREEN]
        
        for i, (key, val) in enumerate(metrics):
            row, col = i // 2, i % 2
            mx = 5.2 + col * 2.0
            my = 0.7 + row * 1.35
            
            # Shadow
            shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(mx + 0.04), Inches(my + 0.04), Inches(1.9), Inches(1.2))
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = self._hex_to_rgb("#C0C0C0")
            shadow.line.fill.background()
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(mx), Inches(my), Inches(1.9), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb(metric_colors[i])
            card.line.fill.background()
            
            # Value - LARGE
            val_txt = slide.shapes.add_textbox(Inches(mx), Inches(my + 0.2), Inches(1.9), Inches(0.55))
            tf = val_txt.text_frame
            p = tf.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            lbl = slide.shapes.add_textbox(Inches(mx), Inches(my + 0.78), Inches(1.9), Inches(0.35))
            tf = lbl.text_frame
            p = tf.paragraphs[0]
            p.text = key.upper()
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb("#E8E8E8")
            p.alignment = PP_ALIGN.CENTER
        
        # ===== RIGHT COLUMN: Products List =====
        prod_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(9.4), Inches(0.7), Inches(3.7), Inches(0.4))
        prod_header.fill.solid()
        prod_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.PINK)
        prod_header.line.fill.background()
        
        prod_txt = slide.shapes.add_textbox(Inches(9.4), Inches(0.74), Inches(3.7), Inches(0.35))
        tf = prod_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Products & Solutions"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Products with colored bars
        products = data.product_portfolio[:8] if data.product_portfolio else ["Product A"]
        bar_colors = [KelpBrandV2.NAVY, KelpBrandV2.TEAL, KelpBrandV2.PURPLE, KelpBrandV2.ORANGE,
                      KelpBrandV2.CHART_GREEN, KelpBrandV2.PINK, KelpBrandV2.NAVY, KelpBrandV2.TEAL]
        
        for i, prod in enumerate(products[:8]):
            py = 1.2 + i * 0.38
            
            # Color bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(9.4), Inches(py), Inches(0.15), Inches(0.32))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._hex_to_rgb(bar_colors[i % len(bar_colors)])
            bar.line.fill.background()
            
            # Product text - NO TRUNCATION
            ptxt = slide.shapes.add_textbox(Inches(9.6), Inches(py + 0.04), Inches(3.4), Inches(0.3))
            tf = ptxt.text_frame
            p = tf.paragraphs[0]
            p.text = prod  # Full text
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
        
        # ===== MIDDLE ROW: Company Snapshot =====
        snapshot_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(0.2), Inches(3.5), Inches(6.3), Inches(2.0))
        snapshot_bg.fill.solid()
        snapshot_bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        snapshot_bg.line.fill.background()
        
        snap_title = slide.shapes.add_textbox(Inches(0.4), Inches(3.6), Inches(5.9), Inches(0.4))
        tf = snap_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Company Snapshot"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Bullets - FULL TEXT with wrapping
        bullets = data.business_bullets[:3] if data.business_bullets else ["Industry leader"]
        for i, bullet in enumerate(bullets):
            btxt = slide.shapes.add_textbox(Inches(0.4), Inches(4.05 + i * 0.48), Inches(5.9), Inches(0.45))
            tf = btxt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"✓ {bullet}"  # Full text, no truncation
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb("#C8E0FF")
        
        # ===== KEY CUSTOMERS =====
        cust_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(6.7), Inches(3.5), Inches(3.1), Inches(0.35))
        cust_header.fill.solid()
        cust_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.TEAL)
        cust_header.line.fill.background()
        
        cust_txt = slide.shapes.add_textbox(Inches(6.7), Inches(3.53), Inches(3.1), Inches(0.3))
        tf = cust_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Customers"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Customer chips - 2 columns
        customers = data.key_customers[:8] if data.key_customers else ["Client A"]
        for i, cust in enumerate(customers[:8]):
            row, col = i // 2, i % 2
            cx = 6.75 + col * 1.55
            cy = 3.95 + row * 0.38
            
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(cx), Inches(cy), Inches(1.48), Inches(0.32))
            chip.fill.solid()
            chip.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.LIGHT_BLUE)
            chip.line.color.rgb = self._hex_to_rgb(KelpBrandV2.TEAL)
            chip.line.width = Pt(0.75)
            
            chip_txt = slide.shapes.add_textbox(Inches(cx + 0.05), Inches(cy + 0.06), Inches(1.38), Inches(0.24))
            tf = chip_txt.text_frame
            p = tf.paragraphs[0]
            p.text = cust[:18]
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
            p.alignment = PP_ALIGN.CENTER
        
        # ===== INDUSTRIES SERVED =====
        ind_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(9.95), Inches(3.5), Inches(3.15), Inches(0.35))
        ind_header.fill.solid()
        ind_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.PURPLE)
        ind_header.line.fill.background()
        
        ind_txt = slide.shapes.add_textbox(Inches(9.95), Inches(3.53), Inches(3.15), Inches(0.3))
        tf = ind_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Industries Served"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Industry chips
        apps = data.applications[:6] if data.applications else ["Industry A"]
        chip_colors = [KelpBrandV2.LIGHT_PURPLE, KelpBrandV2.LIGHT_BLUE, KelpBrandV2.LIGHT_ORANGE,
                       KelpBrandV2.LIGHT_GREEN, KelpBrandV2.LIGHT_YELLOW, KelpBrandV2.LIGHT_PURPLE]
        
        for i, app in enumerate(apps[:6]):
            row, col = i // 2, i % 2
            ax = 10.0 + col * 1.55
            ay = 3.95 + row * 0.45
            
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(ax), Inches(ay), Inches(1.5), Inches(0.38))
            chip.fill.solid()
            chip.fill.fore_color.rgb = self._hex_to_rgb(chip_colors[i % len(chip_colors)])
            chip.line.color.rgb = self._hex_to_rgb(KelpBrandV2.PURPLE)
            chip.line.width = Pt(1)
            
            chip_txt = slide.shapes.add_textbox(Inches(ax + 0.05), Inches(ay + 0.08), Inches(1.4), Inches(0.28))
            tf = chip_txt.text_frame
            p = tf.paragraphs[0]
            p.text = app[:16]
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.PURPLE)
            p.alignment = PP_ALIGN.CENTER
        
        # ===== BOTTOM ROW: Certifications + Second Image =====
        cert_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.2), Inches(5.65), Inches(6.3), Inches(1.15))
        cert_box.fill.solid()
        cert_box.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.LIGHT_YELLOW)
        cert_box.line.color.rgb = self._hex_to_rgb(KelpBrandV2.ORANGE)
        cert_box.line.width = Pt(1.5)
        
        cert_title = slide.shapes.add_textbox(Inches(0.3), Inches(5.7), Inches(6.1), Inches(0.35))
        tf = cert_title.text_frame
        p = tf.paragraphs[0]
        p.text = "🏆 Certifications & Accreditations"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.ORANGE)
        
        certs = data.certifications[:4] if data.certifications else ["ISO 9001"]
        cert_list = slide.shapes.add_textbox(Inches(0.3), Inches(6.05), Inches(6.1), Inches(0.7))
        tf = cert_list.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = " • ".join([c[:35] for c in certs])  # Reasonable length
        p.font.size = Pt(9)
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
        
        # Second image - 16:9 ratio
        if slide2_images and len(slide2_images) >= 2:
            img_path = slide2_images[1].path if hasattr(slide2_images[1], 'path') else slide2_images[1]
            self._add_image_to_slide(slide, img_path, 6.7, 5.55, width=6.4, height=1.3, send_to_back=False)
        
        # Footer
        self._add_footer(slide)
    
    # ========================================================================
    # SLIDE 3: FINANCIAL DEEP-DIVE (Redesigned - Fuller, More Visual)
    # ========================================================================
    
    def _create_slide3_financials(self, slide, data: EnhancedTeaserData) -> None:
        """Create financial slide with dense, data-rich layout"""
        # Clean white background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    KelpBrandV2.SLIDE_WIDTH, KelpBrandV2.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header
        self._add_header_bar(slide, "Financial Performance & Growth Story")
        
        # ===== TOP ROW: 4 Key Metrics Across Full Width =====
        metrics = list(data.key_metrics_headline.items())[:4]
        metric_colors = [KelpBrandV2.NAVY, KelpBrandV2.ORANGE, KelpBrandV2.TEAL, KelpBrandV2.PURPLE]
        
        for i, (key, val) in enumerate(metrics):
            mx = 0.2 + i * 3.28
            
            # Shadow
            shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(mx + 0.03), Inches(0.73), Inches(3.1), Inches(0.85))
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = self._hex_to_rgb("#B0B0B0")
            shadow.line.fill.background()
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(mx), Inches(0.7), Inches(3.1), Inches(0.85))
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb(metric_colors[i])
            card.line.fill.background()
            
            # Value
            val_txt = slide.shapes.add_textbox(Inches(mx), Inches(0.75), Inches(3.1), Inches(0.48))
            tf = val_txt.text_frame
            p = tf.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            lbl = slide.shapes.add_textbox(Inches(mx), Inches(1.25), Inches(3.1), Inches(0.25))
            tf = lbl.text_frame
            p = tf.paragraphs[0]
            p.text = key.upper()
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb("#E0E0E0")
            p.alignment = PP_ALIGN.CENTER
        
        # ===== LEFT: Bar Chart (Full Height) =====
        rev_data = data.revenue_trend if data.revenue_trend else {
            'categories': ['FY21', 'FY22', 'FY23', 'FY24', 'FY25E'],
            'series': [[800, 950, 1100, 1300, 1550], [120, 150, 180, 210, 260], [70, 90, 110, 140, 175]],
            'series_names': ['Revenue (₹Cr)', 'EBITDA (₹Cr)', 'PAT (₹Cr)']
        }
        
        self._add_chart(slide, 0.2, 1.7, 6.3, 2.9, 'column', rev_data,
                       "Revenue, EBITDA & PAT Trend",
                       [KelpBrandV2.CHART_NAVY, KelpBrandV2.CHART_ORANGE, KelpBrandV2.CHART_PURPLE])
        
        # CAGR Badge
        cagr_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(1.9),
                                             Inches(1.2), Inches(0.65))
        cagr_badge.fill.solid()
        cagr_badge.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.PINK)
        cagr_badge.line.fill.background()
        
        cagr_txt = slide.shapes.add_textbox(Inches(5.2), Inches(1.95), Inches(1.2), Inches(0.6))
        tf = cagr_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "15%+\nCAGR"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # ===== RIGHT: Growth Drivers (Full Height Panel) =====
        growth_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               Inches(6.7), Inches(1.7), Inches(6.4), Inches(2.9))
        growth_panel.fill.solid()
        growth_panel.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.LIGHT_GREEN)
        growth_panel.line.color.rgb = self._hex_to_rgb(KelpBrandV2.CHART_GREEN)
        growth_panel.line.width = Pt(2)
        
        growth_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6.2), Inches(0.4))
        tf = growth_title.text_frame
        p = tf.paragraphs[0]
        p.text = "📈 Growth Drivers & Market Opportunity"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.CHART_GREEN)
        
        # Market intelligence badges (if available from web research)
        badge_y = 2.2
        if data.market_size or data.market_cagr:
            # Market Size badge
            if data.market_size:
                mkt_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   Inches(6.8), Inches(badge_y), Inches(3.0), Inches(0.45))
                mkt_badge.fill.solid()
                mkt_badge.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
                mkt_badge.line.fill.background()
                
                mkt_txt = slide.shapes.add_textbox(Inches(6.85), Inches(badge_y + 0.08), Inches(2.9), Inches(0.35))
                tf = mkt_txt.text_frame
                p = tf.paragraphs[0]
                p.text = f"🌍 TAM: {data.market_size}"
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
            
            # CAGR badge
            if data.market_cagr:
                cagr_badge_x = 10.0 if data.market_size else 6.8
                cagr_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                    Inches(cagr_badge_x), Inches(badge_y), Inches(3.0), Inches(0.45))
                cagr_badge.fill.solid()
                cagr_badge.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.ORANGE)
                cagr_badge.line.fill.background()
                
                cagr_txt = slide.shapes.add_textbox(Inches(cagr_badge_x + 0.05), Inches(badge_y + 0.08), Inches(2.9), Inches(0.35))
                tf = cagr_txt.text_frame
                p = tf.paragraphs[0]
                p.text = f"📊 Industry CAGR: {data.market_cagr}"
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
            
            badge_y += 0.55  # Move callouts down
        
        # Growth callouts
        callouts = data.growth_callouts[:4] if data.growth_callouts else [
            "Strong margin expansion through operational efficiency",
            "New capacity to capture growing market demand",
            "Long-term contracts with marquee clients",
            "Geographic expansion into new markets"
        ]
        
        for i, callout in enumerate(callouts[:4]):
            cy = badge_y + 0.1 + i * 0.52
            ctxt = slide.shapes.add_textbox(Inches(6.9), Inches(cy), Inches(6.0), Inches(0.50))
            tf = ctxt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"✓ {callout}"  # Full text, no truncation
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
        
        # ===== BOTTOM ROW: Donut Chart + Line Chart + Sector Image =====
        # Segment Mix - Donut Chart (more modern than pie)
        seg_data = data.segment_mix if data.segment_mix else {"Products": 60, "Services": 25, "Exports": 15}
        self._add_chart(slide, 0.2, 4.7, 4.3, 2.1, 'donut', 
                       {'labels': list(seg_data.keys()), 'values': list(seg_data.values())},
                       "Revenue Mix by Segment",
                       [KelpBrandV2.CHART_NAVY, KelpBrandV2.CHART_ORANGE, KelpBrandV2.CHART_PURPLE, KelpBrandV2.CHART_TEAL])
        
        # Margin Trend - Line Chart with markers (if we have revenue data)
        if data.revenue_trend and 'series' in data.revenue_trend:
            # Create margin trend from revenue data
            rev_data = data.revenue_trend
            if len(rev_data.get('series', [])) >= 2:
                margin_chart_data = {
                    'categories': rev_data.get('categories', ['FY21', 'FY22', 'FY23', 'FY24']),
                    'series': [rev_data['series'][1]],  # EBITDA trend
                    'series_names': ['EBITDA (₹Cr)']
                }
                self._add_chart(slide, 4.6, 4.7, 4.3, 2.1, 'area',
                               margin_chart_data,
                               "EBITDA Trend",
                               [KelpBrandV2.CHART_GREEN])
            else:
                # Fallback to geographic pie
                geo_data = data.geographic_mix if data.geographic_mix else {"Domestic": 65, "Export": 35}
                self._add_chart(slide, 4.6, 4.7, 4.3, 2.1, 'pie',
                               {'labels': list(geo_data.keys()), 'values': list(geo_data.values())},
                               "Geographic Distribution",
                               [KelpBrandV2.CHART_TEAL, KelpBrandV2.CHART_ORANGE, KelpBrandV2.CHART_PURPLE])
        else:
            # Fallback to geographic pie
            geo_data = data.geographic_mix if data.geographic_mix else {"Domestic": 65, "Export": 35}
            self._add_chart(slide, 4.6, 4.7, 4.3, 2.1, 'pie',
                           {'labels': list(geo_data.keys()), 'values': list(geo_data.values())},
                           "Geographic Distribution",
                           [KelpBrandV2.CHART_TEAL, KelpBrandV2.CHART_ORANGE, KelpBrandV2.CHART_PURPLE])
        
        # Sector image (right side) - 16:9 aspect ratio
        slide3_images = self.slide_images.get('slide3', [])
        if slide3_images and len(slide3_images) >= 1:
            img_path = slide3_images[0].path if hasattr(slide3_images[0], 'path') else slide3_images[0]
            self._add_image_to_slide(slide, img_path, 9.0, 4.7, width=4.1, height=2.1, send_to_back=False)
        
        # Footer
        self._add_footer(slide)
    
    # ========================================================================
    # SLIDE 4: INVESTMENT HIGHLIGHTS (Redesigned - More Engaging)
    # ========================================================================
    
    def _create_slide4_investment(self, slide, data: EnhancedTeaserData) -> None:
        """Create visually engaging investment highlights slide with full space utilization"""
        # Clean white background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    KelpBrandV2.SLIDE_WIDTH, KelpBrandV2.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header
        self._add_header_bar(slide, "Investment Highlights")
        
        # ===== LEFT COLUMN: Hero Image + Investment Thesis (5.2" wide) =====
        # Large sector image (16:9 aspect ratio, properly proportioned)
        slide4_images = self.slide_images.get('slide4', [])
        if slide4_images and len(slide4_images) >= 1:
            img_path = slide4_images[0].path if hasattr(slide4_images[0], 'path') else slide4_images[0]
            # Properly sized hero image - 16:9 ratio = 5.0 x 2.81
            self._add_image_to_slide(slide, img_path, 0.2, 0.7, width=5.2, height=2.92, send_to_back=False)
        else:
            # Fallback decorative gradient box
            visual_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.7),
                                               Inches(5.2), Inches(2.92))
            visual_bg.fill.solid()
            visual_bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.LIGHT_BLUE)
            visual_bg.line.color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
            visual_bg.line.width = Pt(2)
        
        # Investment Thesis Box (below image, full height)
        thesis_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(0.2), Inches(3.72), Inches(5.2), Inches(3.08))
        thesis_box.fill.solid()
        thesis_box.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        thesis_box.line.fill.background()
        
        # Thesis header
        thesis_header = slide.shapes.add_textbox(Inches(0.35), Inches(3.85), Inches(4.9), Inches(0.45))
        tf = thesis_header.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 Investment Thesis"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Thesis bullets - FULL TEXT, no truncation
        thesis_points = data.thesis_points if hasattr(data, 'thesis_points') and data.thesis_points else [
            "Strong market position with proven track record",
            "Multiple growth levers across products & geographies",
            "Superior financial profile with healthy margins",
            "Experienced and motivated management team",
            "Clear path to value creation and exit"
        ]
        
        for i, point in enumerate(thesis_points[:5]):
            pt_txt = slide.shapes.add_textbox(Inches(0.35), Inches(4.35 + i * 0.48), Inches(4.9), Inches(0.45))
            tf = pt_txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"✓ {point}"  # Full text, no truncation
            p.font.size = Pt(11)
            p.font.color.rgb = self._hex_to_rgb("#C8E0FF")
        
        # ===== RIGHT COLUMN: 5 Highlight Cards =====
        colors = [KelpBrandV2.ORANGE, KelpBrandV2.PINK, KelpBrandV2.TEAL, 
                 KelpBrandV2.CHART_GREEN, KelpBrandV2.PURPLE]
        bg_colors = [KelpBrandV2.LIGHT_ORANGE, KelpBrandV2.LIGHT_PURPLE, KelpBrandV2.LIGHT_BLUE,
                    KelpBrandV2.LIGHT_GREEN, KelpBrandV2.LIGHT_PURPLE]
        
        highlights = data.investment_highlights[:5] if data.investment_highlights else [
            {"title": "Market Leader", "description": "Strong position with proprietary capabilities and high entry barriers"},
            {"title": "Blue-Chip Clients", "description": "Long-term relationships with Fortune 500 companies, high retention rate"},
            {"title": "Strategic Infrastructure", "description": "Optimal manufacturing footprint for cost advantage and market access"},
            {"title": "Global Presence", "description": "Strong export franchise across Americas, Europe, and APAC regions"},
            {"title": "Financial Strength", "description": "Consistent growth with healthy returns and low leverage profile"}
        ]
        
        for i, hl in enumerate(highlights):
            hy = 0.7 + i * 1.24  # Space out evenly
            
            # Card shadow
            card_shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                  Inches(5.53), Inches(hy + 0.03), Inches(7.55), Inches(1.12))
            card_shadow.fill.solid()
            card_shadow.fill.fore_color.rgb = self._hex_to_rgb("#C0C0C0")
            card_shadow.line.fill.background()
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(5.5), Inches(hy), Inches(7.55), Inches(1.12))
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb(bg_colors[i])
            card.line.color.rgb = self._hex_to_rgb(colors[i])
            card.line.width = Pt(2)
            
            # Number badge
            num_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(hy + 0.25),
                                                Inches(0.6), Inches(0.6))
            num_badge.fill.solid()
            num_badge.fill.fore_color.rgb = self._hex_to_rgb(colors[i])
            num_badge.line.fill.background()
            
            num_txt = slide.shapes.add_textbox(Inches(5.7), Inches(hy + 0.32), Inches(0.6), Inches(0.5))
            tf = num_txt.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
            p.alignment = PP_ALIGN.CENTER
            
            # Title (larger, bolder)
            title = hl.get('title', f'Highlight {i+1}')
            title_box = slide.shapes.add_textbox(Inches(6.45), Inches(hy + 0.12), Inches(6.4), Inches(0.42))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title  # Full text
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
            
            # Description (below title)
            desc = hl.get('description', '')
            if desc:
                desc_box = slide.shapes.add_textbox(Inches(6.45), Inches(hy + 0.56), Inches(6.4), Inches(0.5))
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = desc  # Full text, no truncation
                p.font.size = Pt(10)
                p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.DARK_GRAY)
        
        # Footer
        self._add_footer(slide)
    
    # ========================================================================
    # SLIDE 5: DISCLAIMER
    # ========================================================================
    
    def _create_slide5_disclaimer(self, slide) -> None:
        """Create professional disclaimer slide"""
        # Dark blue background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    KelpBrandV2.SLIDE_WIDTH, KelpBrandV2.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrandV2.NAVY)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Logo
        self._add_kelp_logo(slide, 0.7, 0.7)
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(10), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Important Notice & Disclaimer"
        p.font.size = Pt(32)
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        
        # Disclaimer text
        disclaimer = """Strictly Private & Confidential. This presentation is prepared by Kelp Global exclusively for the intended recipient and may not be reproduced or distributed without prior written consent. This document is for informational purposes only and does not constitute an offer to sell or a solicitation of an offer to buy any securities.

While all information is obtained from sources believed to be reliable, Kelp Global makes no representation or warranty regarding its accuracy or completeness and accepts no liability for any loss arising from its use.

This presentation contains forward-looking statements and financial projections based on current assumptions; actual results may vary materially. All financial data is sourced from company filings and management representations."""
        
        disc_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(11.5), Inches(3.5))
        tf = disc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = disclaimer
        p.font.size = Pt(13)
        p.font.color.rgb = self._hex_to_rgb(KelpBrandV2.WHITE)
        p.line_spacing = 1.5
    
    # ========================================================================
    # MAIN GENERATION
    # ========================================================================
    
    def generate(self, data: EnhancedTeaserData, filename_prefix: str = "") -> Path:
        """Generate complete enhanced teaser presentation"""
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = KelpBrandV2.SLIDE_WIDTH
        self.prs.slide_height = KelpBrandV2.SLIDE_HEIGHT
        
        # Generate codename if not set
        if not data.codename:
            data.codename = random.choice(PROJECT_CODENAMES)
        
        # Use sector_images from data if slide_images not set
        if not self.slide_images and data.sector_images:
            # Distribute images across slides
            images = list(data.sector_images)
            self.slide_images = {
                'slide1': images[0:2] if len(images) > 1 else images[0:1],
                'slide2': images[2:5] if len(images) > 4 else images[1:3],
                'slide3': images[5:7] if len(images) > 6 else images[3:4],
                'slide4': images[7:9] if len(images) > 8 else images[4:6]
            }
        
        blank_layout = self.prs.slide_layouts[6]
        
        # Count images for this presentation
        total_images = sum(len(imgs) for imgs in self.slide_images.values())
        if total_images > 0:
            print(f"  🖼️ Using {total_images} sector images across slides")
        
        print(f"  📊 Creating Slide 1: Cover - Project {data.codename}")
        slide1 = self.prs.slides.add_slide(blank_layout)
        self._create_slide1_cover(slide1, data)
        
        print(f"  📊 Creating Slide 2: Business Overview (4-Quadrant)")
        slide2 = self.prs.slides.add_slide(blank_layout)
        self._create_slide2_business_overview(slide2, data)
        
        print(f"  📊 Creating Slide 3: Financial Deep-Dive (Multi-Chart)")
        slide3 = self.prs.slides.add_slide(blank_layout)
        self._create_slide3_financials(slide3, data)
        
        print(f"  📊 Creating Slide 4: Investment Highlights")
        slide4 = self.prs.slides.add_slide(blank_layout)
        self._create_slide4_investment(slide4, data)
        
        print(f"  📊 Creating Slide 5: Disclaimer")
        slide5 = self.prs.slides.add_slide(blank_layout)
        self._create_slide5_disclaimer(slide5)
        
        # Save
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_prefix = filename_prefix.replace(" ", "_").replace("/", "_")[:25]
        filename = f"Project_{data.codename}_{safe_prefix}_{timestamp}.pptx"
        output_path = self.output_dir / filename
        
        self.prs.save(str(output_path))
        print(f"  ✓ Saved: {output_path.name}")
        
        return output_path


# ============================================================================
# HELPER: Create data from pipeline content
# ============================================================================

def create_enhanced_teaser_data(content: Dict, sector: str, 
                                enriched_metrics: Any = None,
                                basic_info: Dict = None) -> EnhancedTeaserData:
    """Create EnhancedTeaserData from pipeline outputs"""
    data = EnhancedTeaserData()
    data.sector = sector
    data.codename = random.choice(PROJECT_CODENAMES)
    
    # Business bullets
    if 'business_overview' in content:
        data.business_bullets = content['business_overview'][:5]
    
    # Products, customers, etc from basic_info
    if basic_info:
        data.key_customers = basic_info.get('clients', [])[:8]
        data.product_portfolio = basic_info.get('products', [])[:8]
        data.certifications = basic_info.get('certifications', [])[:6]
        data.applications = basic_info.get('industries', [])[:6]
    
    # From enriched metrics
    if enriched_metrics:
        # At a glance
        data.at_a_glance = {}
        if hasattr(enriched_metrics, 'plant_count') and enriched_metrics.plant_count:
            data.at_a_glance["Facilities"] = f"{enriched_metrics.plant_count} Plants"
        if hasattr(enriched_metrics, 'employee_count') and enriched_metrics.employee_count:
            data.at_a_glance["Team"] = f"{enriched_metrics.employee_count:,}"
        if hasattr(enriched_metrics, 'customer_count') and enriched_metrics.customer_count:
            data.at_a_glance["Clients"] = f"{enriched_metrics.customer_count}+"
        if hasattr(enriched_metrics, 'countries_present') and enriched_metrics.countries_present:
            data.at_a_glance["Countries"] = f"{enriched_metrics.countries_present}+"
        
        # Key metrics headline
        data.key_metrics_headline = {}
        if hasattr(enriched_metrics, 'revenue_latest') and enriched_metrics.revenue_latest:
            rev = enriched_metrics.revenue_latest
            data.key_metrics_headline["Revenue"] = f"₹{rev:.0f} Cr"
        if hasattr(enriched_metrics, 'ebitda_margin') and enriched_metrics.ebitda_margin:
            data.key_metrics_headline["EBITDA Margin"] = f"{enriched_metrics.ebitda_margin:.1f}%"
        if hasattr(enriched_metrics, 'pat_margin') and enriched_metrics.pat_margin:
            data.key_metrics_headline["PAT Margin"] = f"{enriched_metrics.pat_margin:.1f}%"
        if hasattr(enriched_metrics, 'revenue_cagr') and enriched_metrics.revenue_cagr:
            data.key_metrics_headline["Revenue CAGR"] = f"{enriched_metrics.revenue_cagr:.0f}%"
        
        # Revenue trend for charts
        if hasattr(enriched_metrics, 'revenue_trend') and enriched_metrics.revenue_trend:
            years = enriched_metrics.revenue_years if hasattr(enriched_metrics, 'revenue_years') else []
            if years:
                data.revenue_trend = {
                    'categories': [f"FY{str(y)[-2:]}" for y in years[-5:]],
                    'series': [
                        enriched_metrics.revenue_trend[-5:],
                        enriched_metrics.ebitda_trend[-5:] if hasattr(enriched_metrics, 'ebitda_trend') and enriched_metrics.ebitda_trend else [],
                        enriched_metrics.pat_trend[-5:] if hasattr(enriched_metrics, 'pat_trend') and enriched_metrics.pat_trend else []
                    ],
                    'series_names': ['Revenue (₹Cr)', 'EBITDA (₹Cr)', 'PAT (₹Cr)']
                }
        
        # Geographic mix
        if hasattr(enriched_metrics, 'export_percentage') and enriched_metrics.export_percentage:
            exp = enriched_metrics.export_percentage
            data.geographic_mix = {"Domestic": 100 - exp, "Export": exp}
        
        # Certifications
        if hasattr(enriched_metrics, 'certifications') and enriched_metrics.certifications:
            data.certifications = enriched_metrics.certifications[:6]
        
        # Clients
        if hasattr(enriched_metrics, 'key_clients') and enriched_metrics.key_clients:
            data.key_customers = enriched_metrics.key_clients[:8]
    
    # Investment highlights from content
    if 'investment_highlights' in content:
        data.investment_highlights = content['investment_highlights'][:5]
    
    # Growth callouts
    if 'growth_highlights' in content:
        data.growth_callouts = content['growth_highlights'][:3]
    
    # Capex plans
    if 'upcoming_facility' in content:
        data.capex_plans = content['upcoming_facility'][:4]
    
    return data


# ============================================================================
# TEST
# ============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("ENHANCED KELP GENERATOR - TEST")
    print("=" * 60)
    
    # Sample data
    test_data = EnhancedTeaserData(
        codename="Apex",
        sector="Manufacturing & Industrials",
        business_bullets=[
            "Leading manufacturer of precision components with 40+ years of expertise",
            "5 state-of-the-art manufacturing facilities across strategic locations",
            "Serving Fortune 500 clients with 95%+ retention rate",
            "Vertically integrated operations from raw material to finished product",
            "Strong R&D focus with 50+ patents and proprietary processes"
        ],
        key_customers=["Daimler", "Volvo", "Mahindra", "Honda", "JCB", "Cummins"],
        at_a_glance={
            "Facilities": "5 Plants",
            "Team": "2,500+",
            "Clients": "150+",
            "Countries": "25+"
        },
        product_portfolio=["Engine Components", "Transmission Parts", "Chassis Systems",
                          "Turbocharger Parts", "Axle Components", "Industrial Products"],
        applications=["Automotive", "Industrial", "Power", "Marine", "Railways", "Agriculture"],
        certifications=["IATF 16949", "ISO 14001", "ISO 45001", "AS9100D"],
        key_metrics_headline={
            "Revenue": "₹2,500 Cr",
            "EBITDA": "18%",
            "PAT": "12%",
            "CAGR": "15%",
            "RoCE": "22%"
        },
        revenue_trend={
            'categories': ['FY21', 'FY22', 'FY23', 'FY24', 'FY25E'],
            'series': [[1800, 2000, 2200, 2400, 2750], [270, 320, 380, 430, 500], [180, 210, 260, 300, 350]],
            'series_names': ['Revenue (₹Cr)', 'EBITDA (₹Cr)', 'PAT (₹Cr)']
        },
        segment_mix={"Products": 55, "Services": 30, "Aftermarket": 15},
        geographic_mix={"Domestic": 60, "Export": 40},
        growth_callouts=[
            "20% margin improvement through automation initiatives",
            "New EV components division launched with ₹200 Cr order book",
            "Long-term supply agreements with 3 new global OEMs"
        ],
        capex_plans=[
            "₹150 Cr expansion capex approved",
            "35% capacity addition by FY26",
            "New EV-focused facility in Gujarat",
            "Expected 25%+ EBITDA margins"
        ],
        investment_highlights=[
            {"title": "Market leader with proprietary forging technology",
             "description": "Ranked #3 in India; 50+ years of metallurgical expertise with high entry barriers"},
            {"title": "Blue-chip OEM relationships with 95%+ retention",
             "description": "Supplying 8 of top 10 global OEMs; embedded in customer product development"},
            {"title": "Strategically located near raw material sources",
             "description": "5 facilities in industrial clusters; optimal logistics and cost structure"},
            {"title": "Diversified revenue with EV upside",
             "description": "35% revenue from non-auto; ₹200 Cr EV order book secured"},
            {"title": "Strong balance sheet with expansion runway",
             "description": "Net debt/EBITDA <1x; ₹150 Cr capex funded through internal accruals"}
        ]
    )
    
    generator = EnhancedKelpGenerator()
    output = generator.generate(test_data, "Manufacturing_Test")
    print(f"\n✅ Generated: {output}")
