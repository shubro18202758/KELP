"""
Enhanced Generative PPT Generator - Rich & Visually Compelling
Creates professional, data-rich presentations with strong visual impact.
Every run produces unique layouts with comprehensive content display.
"""
import random
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import sys

sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from config.settings import KelpBranding, OUTPUT_DIR


# Branding constants
BRANDING = KelpBranding()
DARK_INDIGO = "#2D234B"
LIGHT_BG = "#F8F9FA"
FONT_PRIMARY = "Arial"

# Kelp color palette
KELP_COLORS = {
    'cyan': "#00BFB3",
    'pink': "#E84D8A", 
    'orange': "#FEB95F",
    'purple': "#7B68EE",
    'teal': "#4ECDC4",
    'indigo': "#2D234B",
    'light_purple': "#9B8BC3",
    'gradient_start': "#FF6987",
    'gradient_end': "#FFA55A",
}


@dataclass 
class RichLayoutConfig:
    """Enhanced layout configuration for rich presentations"""
    # Core styling
    primary_accent: str
    secondary_accent: str
    tertiary_accent: str
    
    # Layout types
    slide1_layout: str  # "executive", "magazine", "dashboard", "corporate", "modern"
    slide2_layout: str  # "financials_heavy", "dual_chart", "metrics_wall", "comparison", "infographic"
    slide3_layout: str  # "investment_thesis", "three_column", "numbered", "icon_grid", "story"
    
    # Typography scale
    hero_title_size: int
    section_title_size: int
    body_size: int
    stat_size: int
    caption_size: int
    
    # Visual preferences
    use_gradients: bool
    use_shadows: bool
    chart_style: str
    icon_style: str
    
    # Seed
    seed: int


class RichLayoutGenerator:
    """Generates rich, visually compelling layout configurations"""
    
    SLIDE1_LAYOUTS = ["executive", "magazine", "dashboard", "corporate", "modern"]
    SLIDE2_LAYOUTS = ["financials_heavy", "dual_chart", "metrics_wall", "comparison", "infographic"]
    SLIDE3_LAYOUTS = ["investment_thesis", "three_column", "numbered", "icon_grid", "story"]
    CHART_STYLES = ["modern", "classic", "minimal", "bold"]
    ICON_STYLES = ["circle", "square", "diamond", "pill"]
    
    @classmethod
    def generate(cls, company_name: str = "", force_unique: bool = True) -> RichLayoutConfig:
        if force_unique:
            seed = hash(f"{company_name}_{datetime.now().isoformat()}_{random.random()}")
        else:
            seed = hash(company_name)
        
        rng = random.Random(seed)
        
        colors = list(KELP_COLORS.values())[:5]
        rng.shuffle(colors)
        
        return RichLayoutConfig(
            primary_accent=colors[0],
            secondary_accent=colors[1],
            tertiary_accent=colors[2],
            slide1_layout=rng.choice(cls.SLIDE1_LAYOUTS),
            slide2_layout=rng.choice(cls.SLIDE2_LAYOUTS),
            slide3_layout=rng.choice(cls.SLIDE3_LAYOUTS),
            hero_title_size=rng.randint(36, 44),
            section_title_size=rng.randint(24, 32),
            body_size=rng.randint(12, 14),
            stat_size=rng.randint(28, 36),
            caption_size=rng.randint(9, 11),
            use_gradients=rng.choice([True, True, False]),
            use_shadows=rng.choice([True, False]),
            chart_style=rng.choice(cls.CHART_STYLES),
            icon_style=rng.choice(cls.ICON_STYLES),
            seed=seed
        )


class RichPPTGenerator:
    """
    Creates visually rich, data-dense PowerPoint presentations.
    Designed to look professional and compelling with maximum information display.
    """
    
    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir or OUTPUT_DIR / "rich_presentations"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.layout: Optional[RichLayoutConfig] = None
        self.prs: Optional[Presentation] = None
        
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    
    def _add_gradient_background(self, slide, start_color: str, end_color: str) -> None:
        """Add gradient background to slide"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.gradient()
        background.fill.gradient_angle = 45
        background.fill.gradient_stops[0].color.rgb = self._hex_to_rgb(start_color)
        background.fill.gradient_stops[1].color.rgb = self._hex_to_rgb(end_color)
        background.line.fill.background()
        # Send to back
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def _add_stat_box(self, slide, x: float, y: float, width: float, height: float,
                      value: str, label: str, accent_color: str, style: str = "filled") -> None:
        """Add a prominent stat box"""
        if style == "filled":
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb(accent_color)
            box.line.fill.background()
            text_color = "#FFFFFF"
        else:  # outlined
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
            box.line.color.rgb = self._hex_to_rgb(accent_color)
            box.line.width = Pt(2)
            text_color = DARK_INDIGO
        
        # Value (large)
        val_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + height*0.15),
            Inches(width - 0.3), Inches(height * 0.5)
        )
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(self.layout.stat_size)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(text_color)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Label (smaller)
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + height*0.65),
            Inches(width - 0.2), Inches(height * 0.3)
        )
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(self.layout.caption_size + 1)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(text_color if style == "filled" else accent_color)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    def _add_section_header(self, slide, x: float, y: float, width: float, 
                            title: str, subtitle: str = "") -> None:
        """Add a section header with optional subtitle"""
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.15), Inches(0.6)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y),
            Inches(width - 0.4), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.layout.section_title_size)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.45),
                Inches(width - 0.4), Inches(0.3)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(self.layout.body_size)
            p.font.color.rgb = self._hex_to_rgb("#666666")
            p.font.name = FONT_PRIMARY
    
    def _add_bullet_list(self, slide, x: float, y: float, width: float, height: float,
                         items: List[str], icon_color: str) -> None:
        """Add a styled bullet list with colored icons"""
        current_y = y
        line_height = height / max(len(items), 1)
        
        for i, item in enumerate(items[:8]):  # Max 8 items
            # Bullet circle
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(current_y + 0.05),
                Inches(0.12), Inches(0.12)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self._hex_to_rgb(icon_color)
            bullet.line.fill.background()
            
            # Text
            txt_box = slide.shapes.add_textbox(
                Inches(x + 0.25), Inches(current_y),
                Inches(width - 0.3), Inches(line_height)
            )
            tf = txt_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item[:150]  # Truncate if too long
            p.font.size = Pt(self.layout.body_size)
            p.font.color.rgb = self._hex_to_rgb("#333333")
            p.font.name = FONT_PRIMARY
            
            current_y += line_height
    
    def _add_numbered_list(self, slide, x: float, y: float, width: float, 
                           items: List[str], colors: List[str]) -> None:
        """Add numbered list with colored number badges"""
        current_y = y
        item_height = 0.7
        
        for i, item in enumerate(items[:6]):
            color = colors[i % len(colors)]
            
            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(current_y),
                Inches(0.4), Inches(0.4)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._hex_to_rgb(color)
            badge.line.fill.background()
            
            # Number
            num_box = slide.shapes.add_textbox(
                Inches(x), Inches(current_y + 0.05),
                Inches(0.4), Inches(0.35)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
            
            # Text
            txt_box = slide.shapes.add_textbox(
                Inches(x + 0.55), Inches(current_y + 0.05),
                Inches(width - 0.6), Inches(item_height - 0.1)
            )
            tf = txt_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item[:200]
            p.font.size = Pt(self.layout.body_size + 1)
            p.font.color.rgb = self._hex_to_rgb("#333333")
            p.font.name = FONT_PRIMARY
            
            current_y += item_height
    
    def _add_icon_box(self, slide, x: float, y: float, width: float, height: float,
                      icon_text: str, title: str, description: str, color: str) -> None:
        """Add an info box with icon-like element"""
        # Background card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
        card.line.color.rgb = self._hex_to_rgb("#E0E0E0")
        card.line.width = Pt(1)
        
        # Icon circle
        icon_size = 0.5
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + (width - icon_size)/2), Inches(y + 0.2),
            Inches(icon_size), Inches(icon_size)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = self._hex_to_rgb(color)
        icon.line.fill.background()
        
        # Icon text (emoji or symbol)
        icon_txt = slide.shapes.add_textbox(
            Inches(x + (width - icon_size)/2), Inches(y + 0.25),
            Inches(icon_size), Inches(icon_size)
        )
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.8),
            Inches(width - 0.2), Inches(0.35)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(y + 1.15),
                Inches(width - 0.2), Inches(height - 1.3)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description[:80]
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb("#666666")
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
    
    def _add_chart(self, slide, x: float, y: float, width: float, height: float,
                   chart_type: str, data: Dict, title: str = "") -> None:
        """Add a styled chart"""
        chart_data = CategoryChartData()
        
        if 'categories' in data and 'values' in data:
            chart_data.categories = data['categories']
            chart_data.add_series(data.get('series_name', 'Value'), tuple(data['values']))
        else:
            # Default sample data
            chart_data.categories = ['FY22', 'FY23', 'FY24', 'FY25E']
            chart_data.add_series('Revenue', (100, 120, 145, 175))
        
        xl_type = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'horizontal_bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'pie': XL_CHART_TYPE.PIE,
            'donut': XL_CHART_TYPE.DOUGHNUT,
            'line': XL_CHART_TYPE.LINE_MARKERS,
        }.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        chart = slide.shapes.add_chart(
            xl_type, Inches(x), Inches(y), Inches(width), Inches(height), chart_data
        ).chart
        
        # Style the chart
        if hasattr(chart, 'plots') and len(chart.plots) > 0:
            plot = chart.plots[0]
            if hasattr(plot, 'series') and len(plot.series) > 0:
                series = plot.series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        
        # Add title above chart if specified
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(y - 0.35),
                Inches(width), Inches(0.3)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
            p.font.name = FONT_PRIMARY
    
    def _add_footer(self, slide) -> None:
        """Add Kelp branded footer"""
        slide_width = self.prs.slide_width.inches
        slide_height = self.prs.slide_height.inches
        
        # Footer bar
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(slide_height - 0.45),
            Inches(slide_width), Inches(0.45)
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        footer.line.fill.background()
        
        # Footer text
        footer_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(slide_height - 0.35),
            Inches(slide_width - 1), Inches(0.25)
        )
        tf = footer_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"Strictly Private & Confidential  |  Prepared by Kelp M&A Advisory  |  {datetime.now().strftime('%B %Y')}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 1 LAYOUTS ==========
    
    def _create_slide1_executive(self, slide, content: Dict) -> None:
        """Executive summary layout - professional and comprehensive"""
        slide_w = self.prs.slide_width.inches
        
        # Header bar with gradient effect
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(slide_w), Inches(1.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        header.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.4),
            Inches(slide_w - 1.5), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Investment Opportunity')
        p.font.size = Pt(self.layout.hero_title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        # Sector tag
        sector_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.2),
            Inches(slide_w - 1.5), Inches(0.4)
        )
        tf = sector_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[ {content.get('sector', 'Industry').upper()} SECTOR ]"
        p.font.size = Pt(14)
        p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        p.font.name = FONT_PRIMARY
        
        # Key stats row
        stats = [
            (content.get('financials', {}).get('revenue', 'N/A'), 'REVENUE'),
            (content.get('financials', {}).get('cagr', 'N/A'), 'GROWTH RATE'),
            (content.get('financials', {}).get('ebitda', 'N/A'), 'EBITDA MARGIN'),
            (str(content.get('employees', 'N/A')), 'EMPLOYEES'),
        ]
        
        stat_width = 2.5
        stat_height = 1.2
        start_x = (slide_w - 4*stat_width - 0.6) / 2
        
        for i, (value, label) in enumerate(stats):
            x = start_x + i * (stat_width + 0.2)
            style = "filled" if i % 2 == 0 else "outlined"
            color = self.layout.primary_accent if i % 2 == 0 else self.layout.secondary_accent
            self._add_stat_box(slide, x, 2.1, stat_width, stat_height, value, label, color, style)
        
        # Company description section
        self._add_section_header(slide, 0.75, 3.6, slide_w/2, "Business Overview")
        
        desc_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.2),
            Inches(slide_w/2 - 0.5), Inches(2.3)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        summary = content.get('summary', 'A leading company in its sector.')
        p.text = summary[:600]
        p.font.size = Pt(self.layout.body_size)
        p.font.color.rgb = self._hex_to_rgb("#333333")
        p.font.name = FONT_PRIMARY
        p.line_spacing = 1.4
        
        # Products/Services section
        products = content.get('products', [])[:6]
        if products:
            self._add_section_header(slide, slide_w/2 + 0.5, 3.6, slide_w/2 - 1, "Products & Services")
            self._add_bullet_list(slide, slide_w/2 + 0.5, 4.2, slide_w/2 - 1, 2.3, 
                                  products, self.layout.secondary_accent)
    
    def _create_slide1_magazine(self, slide, content: Dict) -> None:
        """Magazine-style layout with visual impact"""
        slide_w = self.prs.slide_width.inches
        slide_h = self.prs.slide_height.inches
        
        # Left accent panel
        left_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(4.5), Inches(slide_h - 0.45)
        )
        left_panel.fill.solid()
        left_panel.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        left_panel.line.fill.background()
        
        # Sector on left panel
        sector_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(3.5), Inches(0.6)
        )
        tf = sector_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('sector', 'Industry').upper()
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        p.font.name = FONT_PRIMARY
        
        # Key metric on left
        metric_val = content.get('financials', {}).get('revenue', 'N/A')
        metric_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(3.5), Inches(1.2)
        )
        tf = metric_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(metric_val)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        metric_label = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.7),
            Inches(3.5), Inches(0.4)
        )
        tf = metric_label.text_frame
        p = tf.paragraphs[0]
        p.text = "REVENUE"
        p.font.size = Pt(12)
        p.font.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        p.font.name = FONT_PRIMARY
        
        # Additional stats on left
        stats = [
            (content.get('financials', {}).get('cagr', 'N/A'), 'CAGR'),
            (content.get('financials', {}).get('ebitda', 'N/A'), 'EBITDA'),
        ]
        
        for i, (val, lbl) in enumerate(stats):
            y = 4.5 + i * 1.0
            stat_txt = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(1.8), Inches(0.5))
            tf = stat_txt.text_frame
            p = tf.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = FONT_PRIMARY
            
            lbl_txt = slide.shapes.add_textbox(Inches(2.3), Inches(y + 0.1), Inches(1.5), Inches(0.3))
            tf = lbl_txt.text_frame
            p = tf.paragraphs[0]
            p.text = lbl
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb(self.layout.secondary_accent)
            p.font.name = FONT_PRIMARY
        
        # Right content area
        right_x = 5.0
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(right_x), Inches(0.6),
            Inches(slide_w - right_x - 0.5), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get('title', 'Investment Opportunity')
        p.font.size = Pt(self.layout.hero_title_size - 4)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
        p.font.name = FONT_PRIMARY
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(right_x), Inches(2.0),
            Inches(slide_w - right_x - 0.5), Inches(1.8)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get('summary', '')[:400]
        p.font.size = Pt(self.layout.body_size + 1)
        p.font.color.rgb = self._hex_to_rgb("#333333")
        p.font.name = FONT_PRIMARY
        p.line_spacing = 1.5
        
        # Key highlights
        highlights = content.get('highlights', [])[:4]
        if highlights:
            self._add_section_header(slide, right_x, 4.0, slide_w - right_x - 0.5, "Key Highlights")
            self._add_bullet_list(slide, right_x, 4.6, slide_w - right_x - 0.5, 2.0,
                                  highlights, self.layout.primary_accent)
    
    def _create_slide1_dashboard(self, slide, content: Dict) -> None:
        """Dashboard-style with multiple data panels"""
        slide_w = self.prs.slide_width.inches
        
        # Top title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(slide_w), Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        title_bar.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(slide_w - 1), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{content.get('sector', 'Industry')} Investment Opportunity"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        # Stats row
        stats = [
            (content.get('financials', {}).get('revenue', 'N/A'), 'Revenue'),
            (content.get('financials', {}).get('cagr', 'N/A'), 'Growth'),
            (content.get('financials', {}).get('ebitda', 'N/A'), 'EBITDA'),
            (str(content.get('employees', 'N/A')), 'Team Size'),
            (content.get('founded', 'N/A'), 'Established'),
        ]
        
        stat_width = (slide_w - 1.0) / 5 - 0.15
        for i, (val, lbl) in enumerate(stats):
            x = 0.5 + i * (stat_width + 0.15)
            colors = [self.layout.primary_accent, self.layout.secondary_accent, 
                      self.layout.tertiary_accent, self.layout.primary_accent, 
                      self.layout.secondary_accent]
            self._add_stat_box(slide, x, 1.4, stat_width, 1.1, val, lbl, colors[i], "filled")
        
        # Two-column content
        col_width = (slide_w - 1.2) / 2
        
        # Left: Business overview
        self._add_section_header(slide, 0.5, 2.8, col_width, "Business Overview")
        
        desc_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.4),
            Inches(col_width), Inches(1.5)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.get('summary', '')[:350]
        p.font.size = Pt(self.layout.body_size)
        p.font.color.rgb = self._hex_to_rgb("#333333")
        p.font.name = FONT_PRIMARY
        p.line_spacing = 1.4
        
        # Products below description
        products = content.get('products', [])[:4]
        if products:
            self._add_bullet_list(slide, 0.5, 5.0, col_width, 1.5, products, self.layout.tertiary_accent)
        
        # Right: Key highlights
        self._add_section_header(slide, col_width + 0.7, 2.8, col_width, "Investment Highlights")
        
        highlights = content.get('highlights', [])[:5]
        if highlights:
            self._add_numbered_list(slide, col_width + 0.7, 3.4, col_width,
                                    highlights, [self.layout.primary_accent, 
                                                  self.layout.secondary_accent,
                                                  self.layout.tertiary_accent])
    
    # ========== SLIDE 2 LAYOUTS ==========
    
    def _create_slide2_financials_heavy(self, slide, content: Dict) -> None:
        """Financial data-focused layout - NOW USES REAL DATA"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        self._add_section_header(slide, 0.5, 0.4, slide_w - 1, 
                                 "Financial Performance & Metrics",
                                 f"{content.get('sector', 'Industry')} Investment Analysis")
        
        # Use actual chart data from enriched content if available
        actual_chart_data = content.get('chart_data', [])
        if actual_chart_data and len(actual_chart_data) > 0:
            # Use real revenue data
            rev_chart = actual_chart_data[0]
            chart_data = {
                'categories': rev_chart.get('labels', ['FY22', 'FY23', 'FY24', 'FY25']),
                'values': rev_chart.get('values', [100, 118, 142, 170]),
                'series_name': 'Revenue (₹ Cr)'
            }
        else:
            # Fallback to sample data
            chart_data = {
                'categories': ['FY22', 'FY23', 'FY24', 'FY25E'],
                'values': [100, 118, 142, 170],
                'series_name': 'Revenue Growth'
            }
        self._add_chart(slide, 0.5, 1.3, 6.0, 3.5, 'bar', chart_data, "Revenue Trend")
        
        # Key metrics on right - use actual enriched values
        financials = content.get('financials', {})
        metrics = [
            (financials.get('revenue', '₹500 Cr+'), 'Total Revenue', self.layout.primary_accent),
            (financials.get('cagr', '18%+'), 'Revenue CAGR', self.layout.secondary_accent),
            (financials.get('ebitda', '20%+'), 'EBITDA Margin', self.layout.tertiary_accent),
            (financials.get('pat_margin', 'N/A'), 'PAT Margin', self.layout.primary_accent),
        ]
        
        for i, (val, lbl, color) in enumerate(metrics[:4]):
            if val and val != 'N/A':
                self._add_stat_box(slide, 7.0, 1.3 + i*1.05, 2.8, 0.95, val, lbl, color, "filled")
        
        # Additional metrics row - use enriched data
        additional = [
            (financials.get('employees', str(content.get('employees', '1000+'))), 'Employees'),
            (financials.get('plants', content.get('founded', 'N/A')), 'Facilities'),
            (financials.get('customers', str(len(content.get('locations', [])))), 'Clients'),
            (financials.get('credit_rating', 'Rated'), 'Credit Rating'),
        ]
        
        box_width = (slide_w - 1.0) / 4 - 0.15
        for i, (val, lbl) in enumerate(additional):
            x = 0.5 + i * (box_width + 0.15)
            self._add_stat_box(slide, x, 5.0, box_width, 1.1, val, lbl, 
                               self.layout.primary_accent if i%2==0 else self.layout.secondary_accent, "outlined")
    
    def _create_slide2_dual_chart(self, slide, content: Dict) -> None:
        """Two charts side by side with metrics"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        self._add_section_header(slide, 0.5, 0.4, slide_w - 1,
                                 "Performance Metrics & Analysis")
        
        # Bar chart on left - use actual enriched data
        financials = content.get('financials', {})
        
        # Build metrics from actual data
        rev_cagr = financials.get('cagr', '18%').replace('%', '').replace('+', '').strip()
        ebitda_margin = financials.get('ebitda', '20%').replace('%', '').replace('+', '').strip()
        export_pct = financials.get('export_pct', '40%').replace('%', '').replace('+', '').strip()
        
        try:
            rev_val = float(rev_cagr) if rev_cagr and rev_cagr != 'N/A' else 18
        except:
            rev_val = 18
        try:
            ebitda_val = float(ebitda_margin) if ebitda_margin and ebitda_margin != 'N/A' else 22
        except:
            ebitda_val = 22
        try:
            export_val = float(export_pct) if export_pct and export_pct != 'N/A' else 40
        except:
            export_val = 40
        
        chart_data1 = {
            'categories': ['Revenue\nGrowth', 'EBITDA\nMargin', 'Export\nShare', 'Capacity\nUtil.'],
            'values': [rev_val, ebitda_val, export_val, 85],
            'series_name': 'Key Metrics (%)'
        }
        self._add_chart(slide, 0.5, 1.2, 5.5, 2.8, 'bar', chart_data1, "Key Performance Indicators (%)")
        
        # Pie chart on right
        self._add_chart(slide, 6.5, 1.2, 4.5, 2.8, 'pie', {}, "Revenue Mix")
        
        # Highlights below charts - use actual sector KPIs if available
        self._add_section_header(slide, 0.5, 4.3, slide_w - 1, "Financial Highlights")
        
        # Use enriched sector KPIs or content highlights
        sector_kpis = content.get('sector_kpis', [])
        content_highlights = content.get('highlights', [])
        
        if sector_kpis and len(sector_kpis) >= 4:
            highlights = sector_kpis[:4]
        elif content_highlights:
            highlights = content_highlights[:4]
        else:
            highlights = [
                f"Revenue: {financials.get('revenue', 'Strong growth')} with CAGR of {financials.get('cagr', '15%+')}",
                f"EBITDA Margin: {financials.get('ebitda', 'Healthy')} vs industry average",
                f"Team: {financials.get('employees', 'Growing')} employees across locations",
                f"Credit Rating: {financials.get('credit_rating', 'Investment grade')}"
            ]
        
        # Two columns of highlights
        col_width = (slide_w - 1.2) / 2
        for i, h in enumerate(highlights):
            col = i % 2
            row = i // 2
            x = 0.5 + col * (col_width + 0.2)
            y = 4.9 + row * 0.6
            
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05),
                Inches(0.12), Inches(0.12)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self._hex_to_rgb(
                self.layout.primary_accent if col == 0 else self.layout.secondary_accent
            )
            bullet.line.fill.background()
            
            txt = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y), Inches(col_width - 0.3), Inches(0.5))
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = h
            p.font.size = Pt(self.layout.body_size)
            p.font.color.rgb = self._hex_to_rgb("#333333")
            p.font.name = FONT_PRIMARY
    
    def _create_slide2_metrics_wall(self, slide, content: Dict) -> None:
        """Wall of metrics with visual impact"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(slide_w), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        header.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(slide_w - 1), Inches(0.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Metrics & Performance Dashboard"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        # 3x3 grid of metric boxes - USE ACTUAL ENRICHED DATA
        financials = content.get('financials', {})
        
        metrics = [
            (financials.get('revenue', '₹500Cr'), 'REVENUE', '📊'),
            (financials.get('cagr', '18%'), 'CAGR', '📈'),
            (financials.get('ebitda', '22%'), 'EBITDA', '💰'),
            (financials.get('employees', str(content.get('employees', '2000+'))), 'EMPLOYEES', '👥'),
            (financials.get('plants', content.get('founded', '1985')), 'FACILITIES', '🏢'),
            (financials.get('customers', '50+'), 'CLIENTS', '🤝'),
            (financials.get('export_pct', '40%'), 'EXPORTS', '🌍'),
            (financials.get('credit_rating', 'Rated'), 'CREDIT', '📋'),
            (financials.get('certifications', 'ISO 9001')[:12] if financials.get('certifications') else 'ISO 9001', 'CERTIFIED', '✓'),
        ]
        
        cols = 3
        box_width = (slide_w - 1.0) / cols - 0.2
        box_height = 1.5
        
        colors = [self.layout.primary_accent, self.layout.secondary_accent, 
                  self.layout.tertiary_accent]
        
        for i, (val, lbl, icon) in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = 0.5 + col * (box_width + 0.2)
            y = 1.15 + row * (box_height + 0.15)
            
            # Ensure val is not None
            display_val = str(val) if val else 'N/A'
            self._add_icon_box(slide, x, y, box_width, box_height,
                               icon, lbl, display_val, colors[col])
    
    # ========== SLIDE 3 LAYOUTS ==========
    
    def _create_slide3_investment_thesis(self, slide, content: Dict) -> None:
        """Investment thesis focused layout"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(slide_w), Inches(1.0)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        header.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(slide_w - 1), Inches(0.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Thesis & Key Highlights"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        # Highlights with numbers
        highlights = content.get('highlights', [])[:6]
        if not highlights:
            highlights = [
                "Market leader with significant competitive moat",
                "Strong financial performance with consistent growth",
                "Experienced management with proven track record",
                "Strategic positioning for future market expansion"
            ]
        
        colors = [self.layout.primary_accent, self.layout.secondary_accent, self.layout.tertiary_accent]
        self._add_numbered_list(slide, 0.6, 1.3, slide_w - 1.2, highlights, colors)
        
        # Bottom call-to-action box
        cta_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5),
            Inches(slide_w - 1), Inches(1.0)
        )
        cta_box.fill.solid()
        cta_box.fill.fore_color.rgb = self._hex_to_rgb(LIGHT_BG)
        cta_box.line.color.rgb = self._hex_to_rgb(self.layout.primary_accent)
        cta_box.line.width = Pt(2)
        
        cta_txt = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.7),
            Inches(slide_w - 1.4), Inches(0.6)
        )
        tf = cta_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "This opportunity represents a compelling investment case with strong fundamentals, growth potential, and strategic market positioning. Contact Kelp M&A Advisory for detailed due diligence materials."
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb("#444444")
        p.font.name = FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    def _create_slide3_three_column(self, slide, content: Dict) -> None:
        """Three-column layout with icons"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        self._add_section_header(slide, 0.5, 0.4, slide_w - 1,
                                 "Why Invest?", "Key Investment Highlights")
        
        highlights = content.get('highlights', [])[:6]
        
        # Three columns
        col_width = (slide_w - 1.4) / 3
        col_headers = ["Market Position", "Financial Strength", "Growth Potential"]
        icons = ["🎯", "💹", "🚀"]
        
        for i in range(3):
            x = 0.5 + i * (col_width + 0.2)
            
            # Column header with icon
            icon_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + (col_width - 0.6)/2), Inches(1.1),
                Inches(0.6), Inches(0.6)
            )
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = self._hex_to_rgb(
                [self.layout.primary_accent, self.layout.secondary_accent, self.layout.tertiary_accent][i]
            )
            icon_bg.line.fill.background()
            
            header_txt = slide.shapes.add_textbox(
                Inches(x), Inches(1.8),
                Inches(col_width), Inches(0.4)
            )
            tf = header_txt.text_frame
            p = tf.paragraphs[0]
            p.text = col_headers[i]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(DARK_INDIGO)
            p.font.name = FONT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
            
            # Column content
            col_highlights = highlights[i*2:(i+1)*2] if len(highlights) > i*2 else ["Strong performance indicator"]
            self._add_bullet_list(slide, x + 0.1, 2.4, col_width - 0.2, 3.5,
                                  col_highlights, 
                                  [self.layout.primary_accent, self.layout.secondary_accent, self.layout.tertiary_accent][i])
    
    def _create_slide3_icon_grid(self, slide, content: Dict) -> None:
        """Icon grid with key points"""
        slide_w = self.prs.slide_width.inches
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(slide_w), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(DARK_INDIGO)
        header.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(slide_w - 1), Inches(0.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Highlights"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = FONT_PRIMARY
        
        # 2x3 grid of icon boxes
        highlights = content.get('highlights', [])[:6]
        icons = ['📈', '🏆', '🔧', '🌐', '💼', '🎯']
        titles = ['Growth', 'Leadership', 'Operations', 'Presence', 'Portfolio', 'Strategy']
        
        cols = 3
        box_width = (slide_w - 1.0) / cols - 0.15
        box_height = 1.9
        
        colors = [self.layout.primary_accent, self.layout.secondary_accent, self.layout.tertiary_accent]
        
        for i in range(min(6, len(highlights) or 6)):
            col = i % cols
            row = i // cols
            x = 0.5 + col * (box_width + 0.15)
            y = 1.1 + row * (box_height + 0.15)
            
            desc = highlights[i] if i < len(highlights) else "Key competitive advantage"
            self._add_icon_box(slide, x, y, box_width, box_height,
                               icons[i], titles[i], desc, colors[col])
    
    # ========== MAIN GENERATE METHOD ==========
    
    def generate(self, company_name: str, content: Dict, force_unique: bool = True) -> Path:
        """Generate a rich, visually compelling presentation"""
        
        # Generate unique layout
        self.layout = RichLayoutGenerator.generate(company_name, force_unique)
        
        print(f"   🎨 Layout: {self.layout.slide1_layout}/{self.layout.slide2_layout}/{self.layout.slide3_layout}")
        print(f"   🎨 Colors: {self.layout.primary_accent}, {self.layout.secondary_accent}")
        
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Slide 1: Company Overview
        slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        if self.layout.slide1_layout == "executive":
            self._create_slide1_executive(slide1, content)
        elif self.layout.slide1_layout == "magazine":
            self._create_slide1_magazine(slide1, content)
        else:
            self._create_slide1_dashboard(slide1, content)
        
        self._add_footer(slide1)
        
        # Slide 2: Financial Performance
        slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        if self.layout.slide2_layout == "financials_heavy":
            self._create_slide2_financials_heavy(slide2, content)
        elif self.layout.slide2_layout == "dual_chart":
            self._create_slide2_dual_chart(slide2, content)
        else:
            self._create_slide2_metrics_wall(slide2, content)
        
        self._add_footer(slide2)
        
        # Slide 3: Investment Highlights
        slide3 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        if self.layout.slide3_layout == "investment_thesis":
            self._create_slide3_investment_thesis(slide3, content)
        elif self.layout.slide3_layout == "three_column":
            self._create_slide3_three_column(slide3, content)
        else:
            self._create_slide3_icon_grid(slide3, content)
        
        self._add_footer(slide3)
        
        # Save
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = "".join(c for c in company_name if c.isalnum() or c in "_ ").strip()
        filename = f"{safe_name}_Rich_Teaser_{timestamp}.pptx"
        output_path = self.output_dir / filename
        
        self.prs.save(str(output_path))
        print(f"   ✅ Generated: {output_path.name}")
        
        return output_path


def generate_rich_ppt(company_name: str, content: Dict, output_dir: Path = None) -> Path:
    """Main function to generate rich PPT"""
    generator = RichPPTGenerator(output_dir)
    return generator.generate(company_name, content, force_unique=True)


if __name__ == "__main__":
    test_content = {
        "title": "Leading Manufacturing Company",
        "sector": "Manufacturing",
        "summary": "A leading player in precision engineering with strong export focus and sustainable growth trajectory. The company has established itself as a preferred supplier to global OEMs with a diverse product portfolio and world-class manufacturing capabilities spanning multiple facilities.",
        "highlights": [
            "Market leader with 25%+ share in core precision engineering segment",
            "Strong customer relationships spanning 20+ years with blue-chip OEMs",
            "15% revenue CAGR over last 5 years with improving margins",
            "Expanding into high-margin aerospace and defense verticals",
            "ISO 9001, IATF 16949 certified with best-in-class quality metrics",
            "Strategic partnerships with global technology leaders"
        ],
        "financials": {
            "revenue": "₹500 Cr",
            "cagr": "15%",
            "ebitda": "18%",
        },
        "employees": "2,500+",
        "founded": "1985",
        "products": ["Precision Forged Components", "Machined Parts", "Heat-Treated Products", "Sub-Assemblies", "Aerospace Parts"],
        "locations": ["Pune", "Chennai", "Export Hub"],
    }
    
    print("Testing Rich PPT Generator...")
    for i in range(2):
        print(f"\n--- Run {i+1} ---")
        path = generate_rich_ppt("TestCompany", test_content)
        print(f"   Output: {path}")
