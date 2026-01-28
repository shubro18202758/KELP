"""
Kelp Professional Investment Teaser Generator
==============================================

Creates investment-grade M&A presentations matching the exact Kelp format:

Slide 1: Cover - "Project [Codename]" with gradient background
Slide 2: Business Overview - Clients, Metrics Sidebar, Products, Certifications
Slide 3: Key Financial Metrics - Bar charts, Pie charts, World map, Growth callouts
Slide 4: Investment Highlights - Icon-based highlight cards
Slide 5: Disclaimer - Standard confidentiality notice

Reference: Kelp M&A sample presentation format
"""

import random
import math
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import sys

sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from config.settings import OUTPUT_DIR


# ============================================================================
# KELP BRANDING CONSTANTS - Exact match to reference
# ============================================================================

class KelpBrand:
    """Official Kelp brand colors and fonts from reference"""
    # Primary colors
    DARK_BLUE = "#1B365D"       # Main header bar blue
    NAVY = "#0D2137"            # Deep navy for cover
    TEAL = "#00A5B5"            # Accent teal/cyan
    
    # Gradient colors for cover
    GRADIENT_START = "#1B365D"  # Blue
    GRADIENT_MID = "#4A2C6A"    # Purple
    GRADIENT_END = "#E84B8A"    # Pink
    
    # Accent colors
    ORANGE = "#F5A623"          # Orange accent for highlights
    PINK = "#E84B8A"            # Pink gradient accent  
    PURPLE = "#7B68EE"          # Purple accent
    CYAN = "#00BFB3"            # Cyan for icons
    
    # Neutral
    WHITE = "#FFFFFF"
    LIGHT_GRAY = "#F5F5F5"
    DARK_GRAY = "#333333"
    MID_GRAY = "#666666"
    
    # Chart colors
    CHART_BLUE = "#2B5797"
    CHART_ORANGE = "#E98B39"
    CHART_PURPLE = "#7030A0"
    
    # Fonts
    FONT_PRIMARY = "Arial"
    FONT_HEADING = "Arial"
    
    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)


# ============================================================================
# PROJECT CODENAMES - Investment banker style
# ============================================================================

PROJECT_CODENAMES = [
    "Apex", "Titan", "Aurora", "Phoenix", "Atlas", "Orion", "Cosmos", "Nexus",
    "Horizon", "Stellar", "Quantum", "Vertex", "Zenith", "Prime", "Omega",
    "Galaxy", "Summit", "Eclipse", "Frontier", "Catalyst", "Momentum", "Velocity",
    "Everest", "Pinnacle", "Stratos", "Helios", "Neptune", "Mercury", "Jupiter"
]


def generate_codename(company_name: str = "") -> str:
    """Generate a unique project codename"""
    seed = hash(f"{company_name}_{datetime.now().isoformat()}")
    rng = random.Random(seed)
    return rng.choice(PROJECT_CODENAMES)


# ============================================================================
# SECTOR-SPECIFIC ICONS (Unicode symbols for investment presentations)
# ============================================================================

SECTOR_ICONS = {
    "Manufacturing": {
        "main": "🏭",
        "highlights": ["⚙️", "📊", "🌍", "📈", "💼"],
        "metrics": ["🔧", "📦", "👥", "🏆"]
    },
    "Technology": {
        "main": "💻",
        "highlights": ["🚀", "💡", "🔗", "📱", "☁️"],
        "metrics": ["👨‍💻", "🔒", "📊", "🌐"]
    },
    "Pharma": {
        "main": "💊",
        "highlights": ["🔬", "🧪", "📋", "🏥", "🌿"],
        "metrics": ["💉", "📈", "🔍", "✅"]
    },
    "Logistics": {
        "main": "🚛",
        "highlights": ["📦", "🌍", "⏱️", "🛣️", "✈️"],
        "metrics": ["🚢", "📍", "🔄", "📊"]
    },
    "Entertainment": {
        "main": "🎬",
        "highlights": ["🎭", "📺", "🎪", "🌟", "💫"],
        "metrics": ["🎟️", "👥", "📈", "🏆"]
    },
    "Electronics": {
        "main": "📡",
        "highlights": ["🔌", "⚡", "🛰️", "🔧", "📟"],
        "metrics": ["💡", "🔬", "📊", "🌐"]
    }
}


@dataclass
class CompanyTeaserData:
    """Structured data for investment teaser generation"""
    # Core info
    codename: str = ""
    sector: str = "Manufacturing"
    sub_sector: str = ""
    
    # Business Overview (Slide 2)
    business_overview: List[str] = field(default_factory=list)
    key_customers: List[str] = field(default_factory=list)
    at_a_glance: Dict[str, str] = field(default_factory=dict)  # "2 Facilities", "55+ Grades"
    product_portfolio: List[str] = field(default_factory=list)
    applications: List[str] = field(default_factory=list)
    certifications: List[str] = field(default_factory=list)
    
    # Financial Metrics (Slide 3)
    key_metrics_bar: Dict[str, str] = field(default_factory=dict)  # "~19.5% EBITDA", "32.4% RoCE"
    revenue_data: Dict[str, List] = field(default_factory=dict)  # {"years": [], "revenue": [], "ebitda": [], "pat": []}
    revenue_cagr: str = ""
    growth_highlights: List[str] = field(default_factory=list)
    upcoming_facility: List[str] = field(default_factory=list)
    export_domestic_split: Dict[str, float] = field(default_factory=dict)  # {"Export": 72, "Domestic": 28}
    
    # Investment Highlights (Slide 4)
    investment_highlights: List[Dict[str, str]] = field(default_factory=list)  # [{"title": "", "description": ""}]
    
    # Sources for citation
    sources: List[str] = field(default_factory=list)


class KelpProfessionalGenerator:
    """
    Creates investment-grade PowerPoint presentations matching Kelp format.
    
    Output: 5-slide professional investment teaser with:
    - Native editable charts
    - Kelp branding
    - Sector-specific content structure
    - Anonymized company identity
    """
    
    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir or OUTPUT_DIR / "kelp_presentations"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs: Optional[Presentation] = None
        self.brand = KelpBrand()
        
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _add_kelp_logo(self, slide, x: float, y: float, size: float = 0.8) -> None:
        """Add Kelp logo placeholder (stylized K)"""
        # Logo container
        logo_box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(size * 2), Inches(size)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        
        # Kelp text with icon
        run1 = p.add_run()
        run1.text = "✦ "
        run1.font.size = Pt(18)
        run1.font.color.rgb = self._hex_to_rgb(KelpBrand.PINK)
        
        run2 = p.add_run()
        run2.text = "Kelp"
        run2.font.size = Pt(20)
        run2.font.bold = True
        run2.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        run2.font.name = KelpBrand.FONT_HEADING
    
    def _add_kelp_logo_dark(self, slide, x: float, y: float) -> None:
        """Add Kelp logo for light backgrounds"""
        logo_box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(1.5), Inches(0.6)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        
        run1 = p.add_run()
        run1.text = "✦ "
        run1.font.size = Pt(16)
        run1.font.color.rgb = self._hex_to_rgb(KelpBrand.PINK)
        
        run2 = p.add_run()
        run2.text = "Kelp"
        run2.font.size = Pt(18)
        run2.font.bold = True
        run2.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        run2.font.name = KelpBrand.FONT_HEADING
    
    def _add_footer(self, slide) -> None:
        """Add standard Kelp footer"""
        footer = slide.shapes.add_textbox(
            Inches(0), Inches(7.1), Inches(13.333), Inches(0.35)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
        p.font.size = Pt(9)
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.MID_GRAY)
        p.font.name = KelpBrand.FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    def _create_gradient_background(self, slide, style: str = "cover") -> None:
        """Create gradient background matching Kelp style"""
        if style == "cover":
            # Dark gradient with geometric patterns
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                KelpBrand.SLIDE_WIDTH, KelpBrand.SLIDE_HEIGHT
            )
            bg.fill.gradient()
            bg.fill.gradient_angle = 135
            stops = bg.fill.gradient_stops
            stops[0].color.rgb = self._hex_to_rgb(KelpBrand.NAVY)
            stops[1].color.rgb = self._hex_to_rgb(KelpBrand.GRADIENT_MID)
            bg.line.fill.background()
            
            # Accent stripe on left
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(0.4), KelpBrand.SLIDE_HEIGHT
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
            stripe.line.fill.background()
            
            # Geometric overlay shapes (triangles/diamonds)
            self._add_geometric_overlay(slide)
            
            # Send all to back
            for shape in [bg, stripe]:
                spTree = slide.shapes._spTree
                sp = shape._element
                spTree.remove(sp)
                spTree.insert(2, sp)
    
    def _add_geometric_overlay(self, slide) -> None:
        """Add subtle geometric patterns like in reference"""
        # Pink accent triangle top-right
        tri1 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, Inches(8), Inches(-1),
            Inches(6), Inches(5)
        )
        tri1.fill.solid()
        tri1.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.PINK)
        tri1.fill.fore_color.brightness = 0.1
        tri1.line.fill.background()
        
        # Purple overlay
        tri2 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, Inches(6), Inches(2),
            Inches(4), Inches(4)
        )
        tri2.fill.solid()
        tri2.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.PURPLE)
        tri2.fill.fore_color.brightness = 0.2
        tri2.line.fill.background()
    
    # ========================================================================
    # SLIDE 1: COVER SLIDE
    # ========================================================================
    
    def _create_slide1_cover(self, slide, data: CompanyTeaserData) -> None:
        """Create professional cover slide with project codename"""
        # Background
        self._create_gradient_background(slide, "cover")
        
        # Kelp logo top-left
        self._add_kelp_logo(slide, 0.8, 0.6)
        
        # Project Codename (main title)
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.2), Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Project {data.codename}"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        # Subtitle line
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.2), Inches(6), Inches(0.8)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Brief"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        # Website placeholder
        web_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.8), Inches(4), Inches(0.4)
        )
        tf = web_box.text_frame
        p = tf.paragraphs[0]
        p.text = "kelpglobal.com"
        p.font.size = Pt(14)
        p.font.color.rgb = self._hex_to_rgb("#AAAAAA")
        p.font.name = KelpBrand.FONT_PRIMARY
    
    # ========================================================================
    # SLIDE 2: BUSINESS OVERVIEW
    # ========================================================================
    
    def _create_slide2_business_overview(self, slide, data: CompanyTeaserData) -> None:
        """Create business overview slide matching Kelp reference exactly"""
        # White background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, KelpBrand.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        bg.line.fill.background()
        
        # Send to back
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header bar (dark blue)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, Inches(0.6)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        header.line.fill.background()
        
        # Header title
        header_txt = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.12), Inches(4), Inches(0.45)
        )
        tf = header_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Business Overview"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        # Kelp logo in header
        self._add_kelp_logo(slide, 11.5, 0.08, 0.5)
        
        # ---- LEFT COLUMN: Business Overview Box ----
        self._add_section_box(slide, 0.3, 0.8, 5.2, 3.5, 
                             "Business Overview:", data.business_overview,
                             KelpBrand.DARK_BLUE)
        
        # ---- CENTER: Key Select Customers ----
        cust_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(0.8),
            Inches(3.2), Inches(0.45)
        )
        cust_header.fill.solid()
        cust_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.TEAL)
        cust_header.line.fill.background()
        
        cust_txt = slide.shapes.add_textbox(
            Inches(5.7), Inches(0.85), Inches(3.2), Inches(0.4)
        )
        tf = cust_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Select Customers"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        p.alignment = PP_ALIGN.CENTER
        
        # Customer names (simulate logos with styled text)
        self._add_customer_logos(slide, 5.7, 1.4, 3.2, 2.8, data.key_customers)
        
        # ---- RIGHT COLUMN: At a Glance ----
        self._add_at_a_glance(slide, 9.1, 0.8, 4, 3.5, data.at_a_glance, 
                             data.codename, data.certifications)
        
        # ---- BOTTOM: Product Portfolio & Applications ----
        # Section header bar
        prod_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.4),
            Inches(8.2), Inches(0.4)
        )
        prod_header.fill.solid()
        prod_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.TEAL)
        prod_header.line.fill.background()
        
        prod_txt = slide.shapes.add_textbox(
            Inches(0.4), Inches(4.45), Inches(8), Inches(0.35)
        )
        tf = prod_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Growth led through its growing product portfolio and offering solutions to diverse sectors"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Product Portfolio boxes (left)
        self._add_product_grid(slide, 0.5, 4.9, 3.8, 2, data.product_portfolio)
        
        # Applications icons (right)
        self._add_applications_grid(slide, 4.5, 4.9, 4, 2, data.applications)
        
        # Footer
        self._add_footer(slide)
    
    def _add_section_box(self, slide, x: float, y: float, w: float, h: float,
                         title: str, bullets: List[str], border_color: str) -> None:
        """Add business overview section box with bullets"""
        # Border box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        box.line.color.rgb = self._hex_to_rgb(border_color)
        box.line.width = Pt(1.5)
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.1), Inches(y + 0.1),
            Inches(w - 0.2), Inches(0.4)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(border_color)
        title_bar.line.fill.background()
        
        title_txt = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(0.35)
        )
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        # Bullets
        bullet_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.6), Inches(w - 0.3), Inches(h - 0.7)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"■ {bullet}"
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
            p.font.name = KelpBrand.FONT_PRIMARY
            p.space_after = Pt(6)
    
    def _add_customer_logos(self, slide, x: float, y: float, w: float, h: float,
                           customers: List[str]) -> None:
        """Add customer names styled as logo placeholders"""
        # Show top 4 customers prominently
        top_customers = customers[:4] if customers else ["Customer A", "Customer B", "Customer C", "Customer D"]
        
        positions = [
            (x + 0.2, y + 0.1, w - 0.4, 0.6),
            (x + 0.2, y + 0.75, w - 0.4, 0.6),
            (x + 0.2, y + 1.4, w - 0.4, 0.6),
            (x + 0.2, y + 2.05, w - 0.4, 0.6),
        ]
        
        colors = [KelpBrand.DARK_BLUE, KelpBrand.TEAL, KelpBrand.ORANGE, KelpBrand.PURPLE]
        
        for i, (cx, cy, cw, ch) in enumerate(positions[:len(top_customers)]):
            # Customer box
            cbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy),
                Inches(cw), Inches(ch)
            )
            cbox.fill.solid()
            cbox.fill.fore_color.rgb = self._hex_to_rgb("#F8F9FA")
            cbox.line.color.rgb = self._hex_to_rgb("#E0E0E0")
            cbox.line.width = Pt(1)
            
            # Customer name (styled as logo)
            ctxt = slide.shapes.add_textbox(
                Inches(cx + 0.1), Inches(cy + 0.15), Inches(cw - 0.2), Inches(ch - 0.1)
            )
            tf = ctxt.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = top_customers[i]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(colors[i % len(colors)])
            p.font.name = KelpBrand.FONT_HEADING
            p.alignment = PP_ALIGN.CENTER
        
        # "...and others" text
        if len(customers) > 4:
            other_txt = slide.shapes.add_textbox(
                Inches(x), Inches(y + h - 0.4), Inches(w), Inches(0.35)
            )
            tf = other_txt.text_frame
            p = tf.paragraphs[0]
            p.text = f"...and {len(customers) - 4} others"
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.MID_GRAY)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_at_a_glance(self, slide, x: float, y: float, w: float, h: float,
                         metrics: Dict[str, str], codename: str, 
                         certifications: List[str]) -> None:
        """Add 'At a Glance' sidebar with key metrics"""
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(0.45)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.ORANGE)
        header.line.fill.background()
        
        htxt = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.05), Inches(w), Inches(0.4)
        )
        tf = htxt.text_frame
        p = tf.paragraphs[0]
        p.text = f"{codename} at a glance"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Metrics list
        metric_y = y + 0.55
        icons = ["🏭", "📊", "🔬", "👥", "🌍"]
        
        default_metrics = {
            "Facilities": "5+ Plants",
            "Products": "200+ SKUs", 
            "R&D": "In-house Lab",
            "Team": "1000+ Employees"
        }
        
        display_metrics = metrics if metrics else default_metrics
        
        for i, (key, value) in enumerate(list(display_metrics.items())[:4]):
            icon = icons[i % len(icons)]
            
            # Icon
            icon_txt = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(metric_y), Inches(0.4), Inches(0.4)
            )
            tf = icon_txt.text_frame
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(16)
            
            # Value (large)
            val_txt = slide.shapes.add_textbox(
                Inches(x + 0.5), Inches(metric_y - 0.05), Inches(w - 0.6), Inches(0.35)
            )
            tf = val_txt.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
            
            # Label
            lbl_txt = slide.shapes.add_textbox(
                Inches(x + 0.5), Inches(metric_y + 0.28), Inches(w - 0.6), Inches(0.3)
            )
            tf = lbl_txt.text_frame
            p = tf.paragraphs[0]
            p.text = key
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.MID_GRAY)
            
            metric_y += 0.65
        
        # Certifications header
        cert_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(metric_y + 0.1),
            Inches(w), Inches(0.35)
        )
        cert_header.fill.solid()
        cert_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.TEAL)
        cert_header.line.fill.background()
        
        cert_txt = slide.shapes.add_textbox(
            Inches(x), Inches(metric_y + 0.12), Inches(w), Inches(0.32)
        )
        tf = cert_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Certifications"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Certification badges
        cert_y = metric_y + 0.55
        certs = certifications[:6] if certifications else ["ISO 9001", "ISO 14001", "IATF 16949"]
        
        # Create grid of cert badges
        cert_txt = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(cert_y), Inches(w - 0.2), Inches(0.8)
        )
        tf = cert_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = " | ".join(certs[:4])
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_product_grid(self, slide, x: float, y: float, w: float, h: float,
                          products: List[str]) -> None:
        """Add product portfolio as dashed boxes"""
        # Label
        lbl = slide.shapes.add_textbox(
            Inches(x - 0.3), Inches(y + 0.3), Inches(0.3), Inches(1.5)
        )
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = "Product Portfolio"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        # Note: Rotation would need OOXML manipulation
        
        prods = products[:8] if products else ["Product A", "Product B", "Product C", "Product D"]
        
        cols = 2
        rows = min(4, (len(prods) + 1) // 2)
        box_w = (w - 0.3) / cols
        box_h = h / rows
        
        for i, prod in enumerate(prods[:8]):
            row = i // cols
            col = i % cols
            
            bx = x + col * box_w
            by = y + row * box_h
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by),
                Inches(box_w - 0.1), Inches(box_h - 0.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
            box.line.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
            box.line.width = Pt(1)
            box.line.dash_style = 2  # Dashed
            
            txt = slide.shapes.add_textbox(
                Inches(bx + 0.05), Inches(by + 0.1),
                Inches(box_w - 0.2), Inches(box_h - 0.15)
            )
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = prod[:30]
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_applications_grid(self, slide, x: float, y: float, w: float, h: float,
                               applications: List[str]) -> None:
        """Add applications as icon boxes"""
        apps = applications[:6] if applications else ["Industrial", "Automotive", "Power", "Marine"]
        
        app_icons = ["🏭", "🚗", "⚡", "🚢", "🏗️", "🌾", "💊", "🍞"]
        
        cols = 3
        rows = 2
        box_w = w / cols
        box_h = h / rows
        
        for i, app in enumerate(apps[:6]):
            row = i // cols
            col = i % cols
            
            bx = x + col * box_w
            by = y + row * box_h
            
            # Icon
            icon_txt = slide.shapes.add_textbox(
                Inches(bx + (box_w - 0.5) / 2), Inches(by),
                Inches(0.5), Inches(0.5)
            )
            tf = icon_txt.text_frame
            p = tf.paragraphs[0]
            p.text = app_icons[i % len(app_icons)]
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            lbl_txt = slide.shapes.add_textbox(
                Inches(bx), Inches(by + 0.5),
                Inches(box_w), Inches(0.4)
            )
            tf = lbl_txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = app[:15]
            p.font.size = Pt(8)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
            p.alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SLIDE 3: KEY FINANCIAL METRICS AND GROWTH STORY
    # ========================================================================
    
    def _create_slide3_financials(self, slide, data: CompanyTeaserData) -> None:
        """Create financial metrics slide with charts and callouts"""
        # White background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, KelpBrand.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, Inches(0.6)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        header.line.fill.background()
        
        header_txt = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.12), Inches(8), Inches(0.45)
        )
        tf = header_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Financial Metrics and Growth Story"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        self._add_kelp_logo(slide, 11.5, 0.08, 0.5)
        
        # ---- Key Metrics Bar (Top Row) ----
        self._add_metrics_bar(slide, 0.3, 0.75, 8, 0.6, data.key_metrics_bar)
        
        # ---- Main Chart Area ----
        self._add_financial_chart(slide, 0.3, 1.5, 6.5, 3.8, data)
        
        # ---- CAGR Annotation ----
        if data.revenue_cagr:
            cagr_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(2.2), Inches(1.8), Inches(0.5)
            )
            tf = cagr_box.text_frame
            p = tf.paragraphs[0]
            p.text = data.revenue_cagr
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.PINK)
        
        # ---- Right Side: Growth Highlights ----
        self._add_growth_callouts(slide, 7, 1.5, 6, 2.2, data.growth_highlights)
        
        # ---- Upcoming Facility ----
        self._add_upcoming_facility(slide, 9.5, 1.5, 3.5, 2.2, data.upcoming_facility)
        
        # ---- Bottom: Export/Domestic Pie + Map ----
        self._add_export_section(slide, 7, 3.9, 6, 2.8, data.export_domestic_split)
        
        # Footer
        self._add_footer(slide)
    
    def _add_metrics_bar(self, slide, x: float, y: float, w: float, h: float,
                         metrics: Dict[str, str]) -> None:
        """Add horizontal bar of key metrics"""
        # Box
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.TEAL)
        bar.line.fill.background()
        
        # Title inside
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.05), Inches(w), Inches(0.25)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Financial Metrics (FY24)"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        
        # Default metrics
        default_metrics = {
            "EBITDA Margin": "~15%",
            "RoCE": "18%+",
            "RoE": "15%+",
            "Debt Status": "Low Debt"
        }
        display_metrics = metrics if metrics else default_metrics
        
        # Metric boxes
        metric_w = (w - 0.4) / len(display_metrics)
        mx = x + 0.2
        
        colors = [KelpBrand.ORANGE, "#4CAF50", KelpBrand.PURPLE, KelpBrand.PINK]
        
        for i, (key, value) in enumerate(display_metrics.items()):
            # Metric box
            mbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(mx), Inches(y + 0.3),
                Inches(metric_w - 0.1), Inches(0.25)
            )
            mbox.fill.solid()
            mbox.fill.fore_color.rgb = self._hex_to_rgb(colors[i % len(colors)])
            mbox.line.fill.background()
            
            # Value
            vtxt = slide.shapes.add_textbox(
                Inches(mx), Inches(y + 0.32), Inches(metric_w - 0.1), Inches(0.22)
            )
            tf = vtxt.text_frame
            p = tf.paragraphs[0]
            p.text = f"{value} {key}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
            p.alignment = PP_ALIGN.CENTER
            
            mx += metric_w
    
    def _add_financial_chart(self, slide, x: float, y: float, w: float, h: float,
                             data: CompanyTeaserData) -> None:
        """Add main financial bar chart (Revenue/EBITDA/PAT)"""
        chart_data = CategoryChartData()
        
        # Use real data or defaults
        rev_data = data.revenue_data if data.revenue_data else {}
        
        years = rev_data.get('years', ['FY21', 'FY22', 'FY23', 'FY24', 'FY25E'])
        revenue = rev_data.get('revenue', [55, 62, 70, 85, 95])
        ebitda = rev_data.get('ebitda', [8.8, 10, 13.5, 15, 18.5])
        pat = rev_data.get('pat', [5.2, 6, 8.1, 10, 11.2])
        
        chart_data.categories = years
        chart_data.add_series('Revenue (₹Cr)', tuple(revenue))
        chart_data.add_series('EBITDA (₹Cr)', tuple(ebitda))
        chart_data.add_series('PAT (₹Cr)', tuple(pat))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(x), Inches(y), Inches(w), Inches(h),
            chart_data
        )
        
        # Style chart
        chart_obj = chart.chart
        chart_obj.has_legend = True
        chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart_obj.legend.include_in_layout = False
        
        # Color the series
        series_colors = [KelpBrand.DARK_BLUE, KelpBrand.ORANGE, KelpBrand.PURPLE]
        for i, series in enumerate(chart_obj.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self._hex_to_rgb(series_colors[i])
    
    def _add_growth_callouts(self, slide, x: float, y: float, w: float, h: float,
                             highlights: List[str]) -> None:
        """Add growth highlight callout boxes"""
        highlights = highlights[:3] if highlights else [
            "Strong margin expansion driven by product mix improvement",
            "Launched new premium division for higher-value customers",
            "Secured long-term contracts ensuring revenue visibility"
        ]
        
        box_h = h / len(highlights)
        
        for i, highlight in enumerate(highlights):
            by = y + i * box_h
            
            # Callout box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(by),
                Inches(2.4), Inches(box_h - 0.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb("#FFF8E1")  # Light yellow
            box.line.color.rgb = self._hex_to_rgb(KelpBrand.ORANGE)
            box.line.width = Pt(1)
            
            txt = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(by + 0.1),
                Inches(2.2), Inches(box_h - 0.2)
            )
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = highlight[:120]
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
    
    def _add_upcoming_facility(self, slide, x: float, y: float, w: float, h: float,
                               items: List[str]) -> None:
        """Add upcoming facility section"""
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(0.35)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.PINK)
        header.line.fill.background()
        
        htxt = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.3)
        )
        tf = htxt.text_frame
        p = tf.paragraphs[0]
        p.text = "Upcoming Facility"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        
        items = items[:4] if items else [
            "Capex of ~₹50 Cr planned",
            "Adds 30% capacity",
            "Expected revenue: ₹100 Cr+",
            "Superior margins expected"
        ]
        
        # Checkmark items
        item_y = y + 0.45
        for item in items:
            itxt = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(item_y), Inches(w - 0.2), Inches(0.4)
            )
            tf = itxt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"✓ {item[:45]}"
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
            item_y += 0.4
    
    def _add_export_section(self, slide, x: float, y: float, w: float, h: float,
                            split: Dict[str, float]) -> None:
        """Add export/domestic pie chart and global presence"""
        # Global Presence header
        gp_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(0.35)
        )
        gp_header.fill.solid()
        gp_header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.MID_GRAY)
        gp_header.line.fill.background()
        
        gp_txt = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.05), Inches(w), Inches(0.3)
        )
        tf = gp_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Global Presence"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Sales label
        sales_txt = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.5), Inches(1.5), Inches(0.4)
        )
        tf = sales_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Sales"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        
        # Pie chart
        split = split if split else {"Export": 70, "Domestic": 30}
        
        chart_data = CategoryChartData()
        chart_data.categories = list(split.keys())
        chart_data.add_series('Mix', tuple(split.values()))
        
        pie = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(x), Inches(y + 0.9), Inches(2.2), Inches(1.8),
            chart_data
        )
        
        # Style
        chart_obj = pie.chart
        chart_obj.has_legend = True
        chart_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # Color slices
        colors = [KelpBrand.DARK_BLUE, KelpBrand.ORANGE]
        for i, point in enumerate(chart_obj.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self._hex_to_rgb(colors[i % len(colors)])
        
        # World map placeholder (simple text representation)
        map_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 2.5), Inches(y + 0.5),
            Inches(3.3), Inches(2.2)
        )
        map_box.fill.solid()
        map_box.fill.fore_color.rgb = self._hex_to_rgb("#F0F8FF")
        map_box.line.color.rgb = self._hex_to_rgb("#E0E0E0")
        
        map_txt = slide.shapes.add_textbox(
            Inches(x + 2.6), Inches(y + 1.2), Inches(3.1), Inches(1)
        )
        tf = map_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "🌍 Global Footprint\n• Manufacturing: India\n• Clients: 20+ Countries\n• Exports: Americas, EU, APAC"
        p.font.size = Pt(10)
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_GRAY)
    
    # ========================================================================
    # SLIDE 4: INVESTMENT HIGHLIGHTS
    # ========================================================================
    
    def _create_slide4_investment_highlights(self, slide, data: CompanyTeaserData) -> None:
        """Create investment highlights slide with icon-based cards"""
        # White background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, KelpBrand.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, Inches(0.6)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        header.line.fill.background()
        
        header_txt = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.12), Inches(5), Inches(0.45)
        )
        tf = header_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Highlights"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        self._add_kelp_logo(slide, 11.5, 0.08, 0.5)
        
        # Left side: Abstract circular graphic
        self._add_abstract_graphic(slide, 0.5, 1.5, 4.5, 5)
        
        # Right side: Investment highlight cards
        self._add_highlight_cards(slide, 5.5, 0.9, 7.5, 6, data.investment_highlights)
        
        # Footer
        self._add_footer(slide)
    
    def _add_abstract_graphic(self, slide, x: float, y: float, w: float, h: float) -> None:
        """Add abstract circular network graphic"""
        # Outer circle
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.5), Inches(y + 0.5),
            Inches(w - 1), Inches(w - 1)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = self._hex_to_rgb("#E8F4FD")
        outer.line.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        outer.line.width = Pt(3)
        
        # Inner circle
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(y + 1.2),
            Inches(w - 2.4), Inches(w - 2.4)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = self._hex_to_rgb("#B8D9F0")
        inner.line.fill.background()
        
        # Center icon
        center_txt = slide.shapes.add_textbox(
            Inches(x + (w-1)/2), Inches(y + (w-1)/2),
            Inches(1), Inches(1)
        )
        tf = center_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "🌐"
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_highlight_cards(self, slide, x: float, y: float, w: float, h: float,
                             highlights: List[Dict[str, str]]) -> None:
        """Add investment highlight cards with icons"""
        icons = ["🎯", "🏆", "🏭", "🌍", "📈"]
        
        default_highlights = [
            {
                "title": "Leading market position with proprietary capabilities",
                "description": "Ranked among top 5 players globally; offering customized solutions with high entry barriers"
            },
            {
                "title": "Catering to recession-resistant industries",
                "description": "Deep integration with clients' R&D teams has resulted in 90%+ client retention over the last decade"
            },
            {
                "title": "Strategically located manufacturing infrastructure",
                "description": "Minimal logistics cost for raw material procurement; fully automated processing facilities"
            },
            {
                "title": "Established footprint in high-compliance markets",
                "description": "~70% revenue from exports to USA and EU; fully compliant with international standards"
            },
            {
                "title": "Superior financial profile with industry-leading returns",
                "description": "Delivered strong CAGR with sustainable margins; projected continued growth"
            }
        ]
        
        display_highlights = highlights if highlights else default_highlights
        
        card_h = (h - 0.5) / len(display_highlights)
        
        for i, hl in enumerate(display_highlights[:5]):
            cy = y + i * card_h
            icon = icons[i % len(icons)]
            
            title = hl.get('title', f'Investment Highlight {i+1}')
            desc = hl.get('description', 'Key value proposition for investors')
            
            # Icon circle
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(cy + 0.2),
                Inches(0.6), Inches(0.6)
            )
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.TEAL)
            icon_circle.line.fill.background()
            
            # Icon text
            icon_txt = slide.shapes.add_textbox(
                Inches(x + 0.08), Inches(cy + 0.28),
                Inches(0.5), Inches(0.5)
            )
            tf = icon_txt.text_frame
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            
            # Content card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.8), Inches(cy + 0.05),
                Inches(w - 1), Inches(card_h - 0.15)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self._hex_to_rgb("#F8FBFF")
            card.line.color.rgb = self._hex_to_rgb("#D0E4F7")
            card.line.width = Pt(1)
            
            # Title (bold)
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.95), Inches(cy + 0.12),
                Inches(w - 1.3), Inches(0.45)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title[:80]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
            
            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.95), Inches(cy + 0.55),
                Inches(w - 1.3), Inches(card_h - 0.7)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc[:150]
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(KelpBrand.MID_GRAY)
    
    # ========================================================================
    # SLIDE 5: DISCLAIMER
    # ========================================================================
    
    def _create_slide5_disclaimer(self, slide) -> None:
        """Create disclaimer slide"""
        # Dark blue background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            KelpBrand.SLIDE_WIDTH, KelpBrand.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(KelpBrand.DARK_BLUE)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        
        # Kelp logo
        self._add_kelp_logo(slide, 0.8, 0.8, 1)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(10), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Important Notice & Disclaimer"
        p.font.size = Pt(32)
        p.font.bold = False
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_HEADING
        
        # Disclaimer text
        disclaimer = """Strictly Private & Confidential. This presentation is prepared by Kelp Global exclusively for the intended recipient and may not be reproduced or distributed without prior written consent. This document is for informational purposes only and does not constitute an offer to sell or a solicitation of an offer to buy any securities. While all information is obtained from sources believed to be reliable, Kelp Global makes no representation or warranty regarding its accuracy or completeness and accepts no liability for any loss arising from its use. This presentation contains forward-looking statements and financial projections (e.g., FY26-FY28E) based on current assumptions; actual results may vary materially."""
        
        disc_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.2), Inches(11), Inches(3)
        )
        tf = disc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = disclaimer
        p.font.size = Pt(14)
        p.font.color.rgb = self._hex_to_rgb(KelpBrand.WHITE)
        p.font.name = KelpBrand.FONT_PRIMARY
        p.line_spacing = 1.5
    
    # ========================================================================
    # MAIN GENERATION METHOD
    # ========================================================================
    
    def generate(self, data: CompanyTeaserData, filename_prefix: str = "") -> Path:
        """
        Generate complete 5-slide investment teaser presentation.
        
        Args:
            data: CompanyTeaserData with all content
            filename_prefix: Optional prefix for output filename
            
        Returns:
            Path to generated PPTX file
        """
        # Create presentation with 16:9 aspect ratio
        self.prs = Presentation()
        self.prs.slide_width = KelpBrand.SLIDE_WIDTH
        self.prs.slide_height = KelpBrand.SLIDE_HEIGHT
        
        # Generate codename if not provided
        if not data.codename:
            data.codename = generate_codename(filename_prefix)
        
        # Blank layout
        blank_layout = self.prs.slide_layouts[6]
        
        # Create all 5 slides
        print(f"  📊 Creating Slide 1: Cover - Project {data.codename}")
        slide1 = self.prs.slides.add_slide(blank_layout)
        self._create_slide1_cover(slide1, data)
        
        print(f"  📊 Creating Slide 2: Business Overview")
        slide2 = self.prs.slides.add_slide(blank_layout)
        self._create_slide2_business_overview(slide2, data)
        
        print(f"  📊 Creating Slide 3: Key Financial Metrics")
        slide3 = self.prs.slides.add_slide(blank_layout)
        self._create_slide3_financials(slide3, data)
        
        print(f"  📊 Creating Slide 4: Investment Highlights")
        slide4 = self.prs.slides.add_slide(blank_layout)
        self._create_slide4_investment_highlights(slide4, data)
        
        print(f"  📊 Creating Slide 5: Disclaimer")
        slide5 = self.prs.slides.add_slide(blank_layout)
        self._create_slide5_disclaimer(slide5)
        
        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_prefix = filename_prefix.replace(" ", "_").replace("/", "_")[:30]
        filename = f"Project_{data.codename}_{safe_prefix}_{timestamp}.pptx"
        output_path = self.output_dir / filename
        
        # Save
        self.prs.save(str(output_path))
        print(f"  ✓ Saved: {output_path.name}")
        
        return output_path


# ============================================================================
# HELPER: Create teaser data from pipeline content
# ============================================================================

def create_teaser_data_from_content(content: Dict, sector: str, 
                                    enriched_metrics: Any = None) -> CompanyTeaserData:
    """
    Transform pipeline content into CompanyTeaserData structure.
    
    Args:
        content: Generated content from pipeline
        sector: Sector classification
        enriched_metrics: Optional ExtractedMetrics from enrichment engine
        
    Returns:
        CompanyTeaserData ready for PPT generation
    """
    data = CompanyTeaserData()
    data.sector = sector
    data.codename = generate_codename(sector)
    
    # Extract from slide content
    slide1 = content.get('slide1', {})
    slide2 = content.get('slide2', {})
    slide3 = content.get('slide3', {})
    
    # Business Overview (from slide 1)
    desc = slide1.get('company_description', '')
    if desc:
        # Split into bullet points
        sentences = desc.replace('. ', '.|').split('|')
        data.business_overview = [s.strip() for s in sentences if s.strip()][:6]
    
    # Products
    products = content.get('products', []) or slide1.get('key_products', [])
    data.product_portfolio = products[:8] if products else []
    
    # Clients  
    clients = content.get('clients', [])
    data.key_customers = clients[:8] if clients else []
    
    # Certifications
    certs = content.get('certifications', [])
    data.certifications = certs[:6] if certs else []
    
    # Applications / Industries
    industries = content.get('industries', [])
    data.applications = industries[:6] if industries else []
    
    # At a glance metrics
    financials = content.get('financials', {})
    data.at_a_glance = {
        "Facilities": financials.get('plants', '5+ Plants'),
        "Products": f"{len(data.product_portfolio)}+ Products",
        "Team": financials.get('employees', '500+ Employees'),
        "Clients": f"{len(data.key_customers)}+ Clients"
    }
    
    # Use enriched metrics if available
    if enriched_metrics:
        # Financial metrics
        if hasattr(enriched_metrics, 'ebitda_margin') and enriched_metrics.ebitda_margin:
            data.key_metrics_bar["EBITDA Margin"] = f"~{enriched_metrics.ebitda_margin:.1f}%"
        if hasattr(enriched_metrics, 'revenue_latest') and enriched_metrics.revenue_latest:
            rev = enriched_metrics.revenue_latest
            data.key_metrics_bar["Revenue"] = f"₹{rev:.0f} Cr" if rev > 0 else "Growing"
        
        # Revenue trend for chart
        if hasattr(enriched_metrics, 'revenue_trend') and enriched_metrics.revenue_trend:
            years = enriched_metrics.revenue_years if hasattr(enriched_metrics, 'revenue_years') else []
            if years:
                data.revenue_data = {
                    'years': [f"FY{str(y)[-2:]}" for y in years[-5:]],
                    'revenue': enriched_metrics.revenue_trend[-5:],
                    'ebitda': enriched_metrics.ebitda_trend[-5:] if hasattr(enriched_metrics, 'ebitda_trend') else [],
                    'pat': enriched_metrics.pat_trend[-5:] if hasattr(enriched_metrics, 'pat_trend') else []
                }
        
        # CAGR
        if hasattr(enriched_metrics, 'revenue_cagr') and enriched_metrics.revenue_cagr:
            data.revenue_cagr = f"{enriched_metrics.revenue_cagr:.0f}% CAGR(5Y)"
        
        # Export split
        if hasattr(enriched_metrics, 'export_percentage') and enriched_metrics.export_percentage:
            exp = enriched_metrics.export_percentage
            data.export_domestic_split = {"Export": exp, "Domestic": 100 - exp}
        
        # Key clients
        if hasattr(enriched_metrics, 'key_clients') and enriched_metrics.key_clients:
            data.key_customers = enriched_metrics.key_clients[:8]
        
        # At a glance updates
        if hasattr(enriched_metrics, 'plant_count') and enriched_metrics.plant_count:
            data.at_a_glance["Facilities"] = f"{enriched_metrics.plant_count} Plants"
        if hasattr(enriched_metrics, 'employee_count') and enriched_metrics.employee_count:
            data.at_a_glance["Team"] = f"{enriched_metrics.employee_count:,} Employees"
        if hasattr(enriched_metrics, 'certifications') and enriched_metrics.certifications:
            data.certifications = enriched_metrics.certifications[:6]
    
    # Investment highlights
    highlights = slide3.get('investment_highlights', [])
    if highlights:
        data.investment_highlights = []
        for i, h in enumerate(highlights[:5]):
            if isinstance(h, dict):
                data.investment_highlights.append(h)
            else:
                # Split into title and description
                parts = str(h).split(' - ', 1)
                title = parts[0]
                desc = parts[1] if len(parts) > 1 else ""
                data.investment_highlights.append({"title": title, "description": desc})
    
    # Growth highlights
    growth = slide2.get('growth_drivers', [])
    data.growth_highlights = growth[:3] if growth else []
    
    # Upcoming facility
    future = content.get('future_plans', [])
    data.upcoming_facility = future[:4] if future else []
    
    return data


# ============================================================================
# CLI Testing
# ============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("KELP PROFESSIONAL GENERATOR - TEST RUN")
    print("=" * 60)
    
    # Create sample data
    sample_data = CompanyTeaserData(
        codename="Apex",
        sector="Manufacturing",
        business_overview=[
            "Leading manufacturer of natural hydrocolloids and stabilizers",
            "2 state-of-the-art processing units with ~18,000 MTPA capacity",
            "Robust export-oriented model with distribution in Netherlands, USA, Singapore",
            "Sticky vendor relationships with global FMCG giants",
            "Led by second-generation technocrats with deep polymer expertise"
        ],
        key_customers=["Danone", "Halliburton", "Nestlé", "Unilever"],
        at_a_glance={
            "Facilities": "2 Plants",
            "Grades": "55+ Products",
            "R&D": "In-house Lab",
            "Certifications": "FDA, ISO, FSSAI"
        },
        product_portfolio=["Food Grade Guar Gum", "Cassia Gum Powder", "Industrial Thickeners", 
                          "Tamarind Kernel Powder", "Fast Hydration Derivatives", "Blended Stabilizers"],
        applications=["Bakery", "Cosmetics", "Mining", "Textile", "Oil & Gas", "Pet Food"],
        certifications=["FDA", "ISO 9001", "FSSAI", "FSSC 22000", "Sedex", "Kosher"],
        key_metrics_bar={
            "EBITDA Margin": "~19.5%",
            "RoCE": "32.4%",
            "RoE": "21.0%",
            "Debt": "Zero"
        },
        revenue_data={
            'years': ['FY23', 'FY24', 'FY25', 'FY26E', 'FY27E', 'FY28E'],
            'revenue': [55, 62, 70, 95, 115, 138],
            'ebitda': [8.8, 12.1, 13.5, 18.5, 22.4, 26.9],
            'pat': [5.2, 7.5, 8.1, 11.2, 13.8, 16.5]
        },
        revenue_cagr="18% CAGR(5Y)",
        growth_highlights=[
            "Expansion in margins driven by shift towards pharma-grade derivatives",
            "Launched 'Stabilizer Solutions' division for food clients",
            "Contract farming with 5,000+ farmers ensuring supply stability"
        ],
        upcoming_facility=[
            "Capex of ~USD 14mn for Derivatization Plant",
            "Adds 8,000 MT specialized capacity",
            "Expected revenue: ~USD 30mn",
            "Superior EBITDA margins of 22%+"
        ],
        export_domestic_split={"Export": 72, "Domestic": 28},
        investment_highlights=[
            {
                "title": "Leading global player in natural hydrocolloids",
                "description": "Ranked among top 5 exporters of cassia gum globally; offering 50+ customized mesh sizes"
            },
            {
                "title": "Catering to recession-resistant industries",
                "description": "Deep integration with clients' R&D teams has resulted in a 90% client retention rate"
            },
            {
                "title": "Strategically located manufacturing",
                "description": "Located in the heart of the 'Guar Belt' - minimal logistics cost for raw material procurement"
            },
            {
                "title": "Established footprint in high-compliance markets",
                "description": "~70% revenue from exports to USA and EU; fully REACH compliant for European customers"
            },
            {
                "title": "Superior financial profile",
                "description": "Delivered 18% Revenue CAGR; projected to cross USD 130 Mn turnover by FY28"
            }
        ]
    )
    
    # Generate
    generator = KelpProfessionalGenerator()
    output_path = generator.generate(sample_data, "Test_Manufacturing")
    
    print(f"\n✅ Generated: {output_path}")
    print(f"📂 Output directory: {generator.output_dir}")
